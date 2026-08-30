# -*- coding: utf-8 -*-
"""Convert the current PowerPoint into a self-contained, browser-editable deck.

The unopened deck defaults to exact PowerPoint-rendered slide images. Pressing E
switches to native HTML layers where text, shapes and pictures can be edited.
"""

from __future__ import annotations

import base64
import html
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


PPT_PATH = Path(
    "1_Exploration of Dataset/03_Exploration_of_Dataset_"
    "空间域可视化_final_v12.pptx"
)
REFERENCE_DIR = Path(
    "code/_cleanup_archive_20260827/reviews/"
    "review_spatial_ppt_final_v11_html_source"
)
OUTPUT = Path(
    "code/_cleanup_archive_20260827/ppt_intermediates/"
    "空间域可视化_final_v12_generated_editable.html"
)

CANVAS_WIDTH = 1920.0
CANVAS_HEIGHT = 1080.0
DEFAULT_TEXT = "1F2BE0"


def px(value: int | float, scale: float) -> float:
    return round(float(value) * scale, 3)


def css_num(value: float) -> str:
    return f"{value:.3f}".rstrip("0").rstrip(".")


def binary_data_url(data: bytes, mime: str) -> str:
    return f"data:{mime};base64,{base64.b64encode(data).decode('ascii')}"


def path_data_url(path: Path, mime: str = "image/png") -> str:
    return binary_data_url(path.read_bytes(), mime)


def color_hex(color_format, fallback: str | None = None) -> str | None:
    try:
        rgb = color_format.rgb
        if rgb is not None:
            return str(rgb)
    except Exception:
        pass
    return fallback


def fill_opacity(fill) -> float:
    try:
        xfill = fill.fore_color._xFill
        for element in xfill.iter():
            if element.tag.endswith("}alpha"):
                return max(0.0, min(1.0, int(element.get("val")) / 100000.0))
    except Exception:
        pass
    return 1.0


def rgba(hex_color: str, opacity: float) -> str:
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return f"rgba({r},{g},{b},{opacity:.4f})"


def solid_fill_css(fill) -> str:
    try:
        if int(fill.type) != 1:
            return "transparent"
    except Exception:
        return "transparent"
    color = color_hex(fill.fore_color)
    if not color:
        return "transparent"
    return rgba(color, fill_opacity(fill))


def line_css(shape, scale: float) -> tuple[str, str, str]:
    try:
        line_fill = shape.line.fill
        if line_fill.type is None or int(line_fill.type) != 1:
            return "0", "transparent", "solid"
        color = color_hex(shape.line.color, "1F2BE0") or "1F2BE0"
        width_emu = shape.line.width or 9525
        width = max(0.5, px(width_emu, scale))
        dash = "solid"
        try:
            if shape.line.dash_style is not None and "SOLID" not in str(shape.line.dash_style):
                dash = "dashed"
        except Exception:
            pass
        return css_num(width), f"#{color}", dash
    except Exception:
        return "0", "transparent", "solid"


def geometry_name(shape) -> str:
    try:
        geom = shape._element.spPr.prstGeom
        if geom is not None:
            return geom.get("prst") or "rect"
    except Exception:
        pass
    return "rect"


def border_radius(shape, width: float, height: float) -> str:
    geom = geometry_name(shape)
    if geom == "ellipse":
        return "50%"
    if geom in {"roundRect", "round1Rect", "round2SameRect", "round2DiagRect"}:
        return f"{css_num(min(width, height) * 0.16)}px"
    return "0"


def shadow_css(shape) -> str:
    try:
        if "outerShdw" in shape._element.xml:
            return "2px 3px 8px rgba(0,0,0,.20)"
    except Exception:
        pass
    return "none"


def base_geometry(shape, scale: float) -> tuple[float, float, float, float]:
    return (
        px(shape.left, scale),
        px(shape.top, scale),
        px(shape.width, scale),
        px(shape.height, scale),
    )


def text_color(run) -> str:
    try:
        color = color_hex(run.font.color)
        if color:
            return color
    except Exception:
        pass
    return DEFAULT_TEXT


def run_html(run) -> str:
    font = run.font
    name = font.name or "Segoe UI"
    size = float(font.size.pt) * 2.0 if font.size is not None else 28.0
    weight = "700" if font.bold else "400"
    style = "italic" if font.italic else "normal"
    decoration = "underline" if font.underline else "none"
    content = html.escape(run.text).replace("\x0b", "<br>").replace("\n", "<br>")
    return (
        f'<span style="font-family:{html.escape(name)};font-size:{css_num(size)}px;'
        f'font-weight:{weight};font-style:{style};text-decoration:{decoration};'
        f'color:#{text_color(run)}">{content}</span>'
    )


def paragraph_html(paragraph) -> str:
    align = "left"
    if paragraph.alignment is not None:
        label = str(paragraph.alignment)
        if "CENTER" in label:
            align = "center"
        elif "RIGHT" in label:
            align = "right"
        elif "JUSTIFY" in label:
            align = "justify"
    line_height = "1.16"
    if isinstance(paragraph.line_spacing, float):
        line_height = css_num(paragraph.line_spacing)
    content = "".join(run_html(run) for run in paragraph.runs)
    if not content:
        content = html.escape(paragraph.text).replace("\x0b", "<br>").replace("\n", "<br>")
    return (
        f'<div class="paragraph" style="text-align:{align};line-height:{line_height};'
        f'margin:0;padding:0">{content or "&nbsp;"}</div>'
    )


def vertical_alignment(text_frame) -> str:
    label = str(text_frame.vertical_anchor)
    if "BOTTOM" in label:
        return "flex-end"
    if "MIDDLE" in label:
        return "center"
    return "flex-start"


def shape_style(shape, scale: float, *, include_fill: bool = True) -> str:
    left, top, width, height = base_geometry(shape, scale)
    line_width, line_color, line_style = line_css(shape, scale)
    fill = solid_fill_css(shape.fill) if include_fill else "transparent"
    rotation = float(shape.rotation or 0)
    style = [
        "position:absolute",
        f"left:{css_num(left)}px",
        f"top:{css_num(top)}px",
        f"width:{css_num(width)}px",
        f"height:{css_num(height)}px",
        f"background:{fill}",
        f"border:{line_width}px {line_style} {line_color}",
        f"border-radius:{border_radius(shape, width, height)}",
        f"box-shadow:{shadow_css(shape)}",
        "box-sizing:border-box",
        "overflow:hidden",
    ]
    if rotation:
        style.extend([f"transform:rotate({rotation:g}deg)", "transform-origin:center center"])
    return ";".join(style)


def shape_name(shape, index: int) -> str:
    try:
        return shape.name or f"shape-{index}"
    except Exception:
        return f"shape-{index}"


def line_html(shape, index: int, scale: float) -> str:
    left, top, width, height = base_geometry(shape, scale)
    line_width, line_color, line_style = line_css(shape, scale)
    xfrm = shape._element.spPr.xfrm
    flip_h = xfrm.get("flipH") in {"1", "true"}
    flip_v = xfrm.get("flipV") in {"1", "true"}
    x1, x2 = (width, 0) if flip_h else (0, width)
    y1, y2 = (height, 0) if flip_v else (0, height)
    xml = shape._element.xml
    has_arrow = "tailEnd" in xml or "headEnd" in xml
    dash = "6 5" if line_style == "dashed" else "none"
    marker = ' marker-end="url(#arrowhead)"' if has_arrow else ""
    name = html.escape(shape_name(shape, index))
    return (
        f'<svg class="editable-shape line-element" data-shape="{name}" '
        f'style="position:absolute;left:{css_num(left)}px;top:{css_num(top)}px;'
        f'width:{css_num(max(width, 1))}px;height:{css_num(max(height, 1))}px;overflow:visible" '
        f'viewBox="0 0 {css_num(max(width, 1))} {css_num(max(height, 1))}">'
        f'<line x1="{css_num(x1)}" y1="{css_num(y1)}" x2="{css_num(x2)}" y2="{css_num(y2)}" '
        f'stroke="{line_color}" stroke-width="{line_width}" stroke-dasharray="{dash}"{marker}/></svg>'
    )


def picture_html(shape, index: int, scale: float) -> str:
    left, top, width, height = base_geometry(shape, scale)
    source = binary_data_url(shape.image.blob, shape.image.content_type)
    name = html.escape(shape_name(shape, index))
    return (
        f'<img class="editable-shape picture-element" data-shape="{name}" '
        f'src="{source}" alt="{name}" draggable="false" '
        f'style="position:absolute;left:{css_num(left)}px;top:{css_num(top)}px;'
        f'width:{css_num(width)}px;height:{css_num(height)}px;object-fit:fill;overflow:hidden">'
    )


def regular_shape_html(shape, index: int, scale: float) -> str:
    name = html.escape(shape_name(shape, index))
    text = ""
    has_text = bool(getattr(shape, "has_text_frame", False) and shape.text_frame is not None)
    style = shape_style(shape, scale)
    classes = ["editable-shape", "ppt-shape"]
    attributes = [f'data-shape="{name}"', f'style="{style}"']
    if has_text and shape.text.strip():
        tf = shape.text_frame
        left_margin = px(tf.margin_left or 0, scale)
        right_margin = px(tf.margin_right or 0, scale)
        top_margin = px(tf.margin_top or 0, scale)
        bottom_margin = px(tf.margin_bottom or 0, scale)
        content = "".join(paragraph_html(paragraph) for paragraph in tf.paragraphs)
        text = (
            f'<div class="text-content" spellcheck="false" contenteditable="false" '
            f'style="position:absolute;inset:0;display:flex;flex-direction:column;'
            f'justify-content:{vertical_alignment(tf)};padding:{css_num(top_margin)}px '
            f'{css_num(right_margin)}px {css_num(bottom_margin)}px {css_num(left_margin)}px;'
            f'box-sizing:border-box;white-space:pre-wrap;overflow:hidden">{content}</div>'
        )
        classes.append("text-element")
    return f'<div class="{" ".join(classes)}" {" ".join(attributes)}>{text}</div>'


def render_shape(shape, index: int, scale: float) -> str:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return picture_html(shape, index, scale)
    if shape.shape_type == MSO_SHAPE_TYPE.LINE:
        return line_html(shape, index, scale)
    return regular_shape_html(shape, index, scale)


def build_slides(presentation: Presentation) -> str:
    width_scale = CANVAS_WIDTH / presentation.slide_width
    height_scale = CANVAS_HEIGHT / presentation.slide_height
    if abs(width_scale - height_scale) / width_scale > 0.001:
        raise ValueError("Slide scaling is not uniform")
    scale = width_scale
    sections = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        reference_path = REFERENCE_DIR / f"slide-{slide_number}.png"
        if not reference_path.is_file():
            raise FileNotFoundError(reference_path)
        reference = path_data_url(reference_path)
        native = "\n".join(
            render_shape(shape, index, scale) for index, shape in enumerate(slide.shapes)
        )
        sections.append(
            f'''<section class="slide" data-slide="{slide_number}" aria-label="Slide {slide_number}">
  <svg class="marker-defs" width="0" height="0" aria-hidden="true"><defs>
    <marker id="arrowhead" markerWidth="8" markerHeight="8" refX="7" refY="3" orient="auto" markerUnits="strokeWidth">
      <path d="M0,0 L0,6 L8,3 z" fill="#1F2BE0"></path>
    </marker>
  </defs></svg>
  <div class="native-layer">{native}</div>
  <img class="slide-reference" src="{reference}" alt="PowerPoint reference for slide {slide_number}" draggable="false">
</section>'''
        )
    return "\n".join(sections)


TEMPLATE = r'''<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1, viewport-fit=cover">
  <meta name="generator" content="Editable HTML reconstructed from PowerPoint">
  <title>Editable spatial-domain CO₂ visualisation</title>
  <style>
    :root { color-scheme: dark; }
    * { box-sizing: border-box; }
    html, body { width: 100%; height: 100%; margin: 0; overflow: hidden; }
    body { display:grid; place-items:center; background:#111; font-family:"Segoe UI",Arial,sans-serif; user-select:none; }
    .stage { position:relative; width:min(100vw,calc(100vh * 16 / 9)); height:min(100vh,calc(100vw * 9 / 16)); background:#fff; overflow:hidden; box-shadow:0 0 38px rgba(0,0,0,.55); }
    .slide { position:absolute; inset:0; display:none; width:1920px; height:1080px; background:#fff; transform-origin:top left; overflow:hidden; }
    .slide.active { display:block; }
    .native-layer { position:absolute; inset:0; width:1920px; height:1080px; visibility:hidden; }
    .slide-reference { position:absolute; inset:0; z-index:10000; display:block; width:1920px; height:1080px; object-fit:fill; pointer-events:none; }
    body[data-edited="true"] .native-layer, body.editing .native-layer { visibility:visible; }
    body[data-edited="true"] .slide-reference, body.editing .slide-reference { display:none; }
    body.editing.reference-on .slide-reference { display:block; opacity:.32; }
    .marker-defs { position:absolute; }
    .editable-shape { pointer-events:none; }
    body.editing .editable-shape { pointer-events:auto; }
    body.editing .editable-shape:hover { outline:1px dashed rgba(31,43,224,.55); outline-offset:1px; }
    body.editing .editable-shape.selected { outline:3px solid #1F2BE0; outline-offset:2px; resize:both; }
    body.editing .picture-element { cursor:default; }
    body.editing .text-content { cursor:text; user-select:text; }
    body.editing.alt-down .editable-shape { cursor:move; }
    .hit { position:absolute; inset-block:0; z-index:11000; width:18%; cursor:pointer; }
    .hit.prev { left:0; } .hit.next { right:0; }
    body.editing .hit { display:none; }
    .toolbar { position:absolute; left:50%; bottom:18px; z-index:12000; display:flex; align-items:center; gap:8px; transform:translateX(-50%); padding:8px 12px; border:1px solid rgba(255,255,255,.24); border-radius:999px; background:rgba(0,0,0,.76); color:#fff; opacity:0; transition:opacity .18s ease; backdrop-filter:blur(8px); }
    body.show-controls .toolbar, body.editing .toolbar, .toolbar:focus-within { opacity:1; }
    .toolbar button { height:31px; min-width:34px; padding:0 11px; border:0; border-radius:999px; background:rgba(255,255,255,.14); color:#fff; font:600 13px/1 "Segoe UI",Arial,sans-serif; cursor:pointer; }
    .toolbar button:hover, .toolbar button.active { background:#1F2BE0; }
    .counter { min-width:58px; text-align:center; font-size:13px; }
    .edit-help { position:absolute; top:14px; left:50%; z-index:12000; display:none; transform:translateX(-50%); padding:8px 14px; border-radius:999px; background:rgba(31,43,224,.91); color:#fff; font-size:13px; }
    body.editing .edit-help { display:block; }
    @media print {
      @page { size:13.333in 7.5in; margin:0; }
      html, body { display:block; width:auto; height:auto; overflow:visible; background:#fff; }
      .stage { width:1920px!important; height:auto!important; box-shadow:none; overflow:visible; }
      .slide { position:relative; display:block; transform:none!important; page-break-after:always; }
      .native-layer { visibility:visible!important; }
      .slide-reference { display:none!important; }
      .toolbar, .edit-help, .hit { display:none!important; }
    }
  </style>
</head>
<body data-edited="false">
  <main class="stage" id="stage" aria-label="Editable spatial-domain CO₂ visualisation">
    __SLIDES__
    <div class="hit prev" aria-label="Previous slide"></div>
    <div class="hit next" aria-label="Next slide"></div>
    <div class="edit-help">Edit mode · click text to type · click an element to resize · Alt-drag to move · double-click an image to replace</div>
    <nav class="toolbar" aria-label="Presentation and editing controls">
      <button type="button" data-action="previous" title="Previous slide">←</button>
      <span class="counter" aria-live="polite"></span>
      <button type="button" data-action="next" title="Next slide">→</button>
      <button type="button" data-action="edit" title="Toggle edit mode (E)">Edit</button>
      <button type="button" data-action="reference" title="Toggle PowerPoint reference (R)">Reference</button>
      <button type="button" data-action="save" title="Download edited HTML">Save HTML</button>
      <button type="button" data-action="fullscreen" title="Toggle fullscreen (F)">⛶</button>
    </nav>
    <input id="image-picker" type="file" accept="image/*" hidden>
  </main>
  <script>
    const body = document.body;
    const stage = document.getElementById('stage');
    const slides = [...document.querySelectorAll('.slide')];
    const counter = document.querySelector('.counter');
    const toolbar = document.querySelector('.toolbar');
    const picker = document.getElementById('image-picker');
    let current = Math.max(0, Math.min(slides.length - 1, Number(new URLSearchParams(location.hash.slice(1)).get('slide') || 1) - 1));
    let selected = null;
    let imageTarget = null;
    let drag = null;
    let controlsTimer;

    function scaleSlides() {
      const scale = Math.min(stage.clientWidth / 1920, stage.clientHeight / 1080);
      slides.forEach(slide => slide.style.transform = `scale(${scale})`);
    }
    function show(index) {
      current = (index + slides.length) % slides.length;
      slides.forEach((slide, i) => slide.classList.toggle('active', i === current));
      counter.textContent = `${current + 1} / ${slides.length}`;
      history.replaceState(null, '', `#slide=${current + 1}`);
      select(null);
    }
    function select(element) {
      selected?.classList.remove('selected');
      selected = element;
      selected?.classList.add('selected');
      if (selected && editing()) markEdited();
    }
    function markEdited() { body.dataset.edited = 'true'; }
    function editing() { return body.classList.contains('editing'); }
    function toggleEdit(force) {
      const enabled = force ?? !editing();
      body.classList.toggle('editing', enabled);
      document.querySelector('[data-action="edit"]').classList.toggle('active', enabled);
      document.querySelectorAll('.text-content').forEach(element => element.contentEditable = enabled ? 'true' : 'false');
      if (!enabled) { body.classList.remove('reference-on'); select(null); }
    }
    function toggleReference() {
      if (!editing()) toggleEdit(true);
      body.classList.toggle('reference-on');
      document.querySelector('[data-action="reference"]').classList.toggle('active', body.classList.contains('reference-on'));
    }
    async function fullscreen() {
      if (!document.fullscreenElement) await stage.requestFullscreen();
      else await document.exitFullscreen();
    }
    function saveHtml() {
      const wasEditing = editing();
      toggleEdit(false);
      body.classList.remove('show-controls', 'reference-on', 'alt-down');
      const clone = document.documentElement.cloneNode(true);
      clone.querySelectorAll('.selected').forEach(element => element.classList.remove('selected'));
      const source = '<!doctype html>\n' + clone.outerHTML;
      const url = URL.createObjectURL(new Blob([source], {type:'text/html;charset=utf-8'}));
      const link = document.createElement('a');
      link.href = url;
      link.download = '空间域可视化_final_v11_edited.html';
      link.click();
      setTimeout(() => URL.revokeObjectURL(url), 1000);
      if (wasEditing) toggleEdit(true);
    }
    function revealControls() {
      body.classList.add('show-controls');
      clearTimeout(controlsTimer);
      controlsTimer = setTimeout(() => { if (!editing()) body.classList.remove('show-controls'); }, 1400);
    }

    document.querySelector('.hit.prev').addEventListener('click', () => show(current - 1));
    document.querySelector('.hit.next').addEventListener('click', () => show(current + 1));
    toolbar.addEventListener('click', event => {
      const action = event.target.closest('button')?.dataset.action;
      if (action === 'previous') show(current - 1);
      if (action === 'next') show(current + 1);
      if (action === 'edit') toggleEdit();
      if (action === 'reference') toggleReference();
      if (action === 'save') saveHtml();
      if (action === 'fullscreen') fullscreen();
    });

    stage.addEventListener('click', event => {
      if (!editing()) return;
      const shape = event.target.closest('.editable-shape');
      if (shape) select(shape);
    });
    stage.addEventListener('input', event => {
      if (event.target.closest('.text-content')) markEdited();
    });
    stage.addEventListener('dblclick', event => {
      if (!editing()) return;
      const image = event.target.closest('.picture-element');
      if (!image) return;
      imageTarget = image;
      picker.click();
    });
    picker.addEventListener('change', () => {
      const file = picker.files?.[0];
      if (!file || !imageTarget) return;
      const reader = new FileReader();
      reader.onload = () => { imageTarget.src = reader.result; markEdited(); };
      reader.readAsDataURL(file);
      picker.value = '';
    });

    stage.addEventListener('pointerdown', event => {
      if (!editing() || !event.altKey) return;
      const shape = event.target.closest('.editable-shape');
      if (!shape) return;
      event.preventDefault();
      select(shape);
      const left = parseFloat(shape.style.left || 0);
      const top = parseFloat(shape.style.top || 0);
      drag = {shape, x:event.clientX, y:event.clientY, left, top};
      shape.setPointerCapture(event.pointerId);
    });
    stage.addEventListener('pointermove', event => {
      if (!drag) return;
      const scale = stage.clientWidth / 1920;
      drag.shape.style.left = `${drag.left + (event.clientX - drag.x) / scale}px`;
      drag.shape.style.top = `${drag.top + (event.clientY - drag.y) / scale}px`;
      markEdited();
    });
    stage.addEventListener('pointerup', () => { drag = null; });
    document.addEventListener('keydown', event => {
      if (event.key === 'Alt') body.classList.add('alt-down');
      const typing = event.target.closest?.('[contenteditable="true"]');
      if (typing && !event.ctrlKey && !event.metaKey) return;
      if (event.key.toLowerCase() === 'e') { event.preventDefault(); toggleEdit(); }
      else if (event.key.toLowerCase() === 'r') { event.preventDefault(); toggleReference(); }
      else if (event.key.toLowerCase() === 'f') fullscreen();
      else if (!editing() && ['ArrowRight','PageDown','Enter',' '].includes(event.key)) { event.preventDefault(); show(current + 1); }
      else if (!editing() && ['ArrowLeft','PageUp','Backspace'].includes(event.key)) { event.preventDefault(); show(current - 1); }
    });
    document.addEventListener('keyup', event => { if (event.key === 'Alt') body.classList.remove('alt-down'); });
    document.addEventListener('mousemove', revealControls);
    document.addEventListener('touchstart', revealControls, {passive:true});
    window.addEventListener('resize', scaleSlides);
    scaleSlides();
    show(current);
  </script>
</body>
</html>
'''


def main() -> None:
    if not PPT_PATH.is_file():
        raise FileNotFoundError(PPT_PATH)
    presentation = Presentation(PPT_PATH)
    if len(presentation.slides) != 6:
        raise ValueError(f"Expected 6 slides, found {len(presentation.slides)}")
    slides = build_slides(presentation)
    document = TEMPLATE.replace("__SLIDES__", slides)
    OUTPUT.write_text(document, encoding="utf-8")
    print(f"{OUTPUT} ({OUTPUT.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
