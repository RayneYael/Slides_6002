# -*- coding: utf-8 -*-
"""
5-slide spatial PPT (16:9) — English narrative.

Each slide has a single BIG TITLE that IS the main point of that slide
(not "this slide shows X").  The image is on the left, a frosted-glass
content card is on the right.  Header: green frosted-glass with point
label.
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu
from lxml import etree

# ──────────────────────────────────────────────────────────────────────────────
# Theme
# ──────────────────────────────────────────────────────────────────────────────
INK = RGBColor(0x1F, 0x2B, 0xE0)          # Cobalt Grid electric ink
INK_SOFT = RGBColor(0x55, 0x60, 0xE5)
MUTED = RGBColor(0x67, 0x6E, 0xA8)
PAGE = RGBColor(0xF0, 0xEB, 0xDE)         # warm cream paper
PAGE_2 = RGBColor(0xE6, 0xE0, 0xCE)
GRID = RGBColor(0xD7, 0xD8, 0xF2)         # ~10% cobalt on paper
ACCENT = INK
ACCENT_DARK = INK
ACCENT_DEEP = RGBColor(0x16, 0x20, 0xA8)
ACCENT_SOFT = GRID
AMBER = INK_SOFT

FONT_HEAD = "Georgia"                     # Newsreader-compatible fallback
FONT_BODY = "Segoe UI"                    # Hanken Grotesk-compatible fallback
FONT_MONO = "Consolas"                    # DM Mono-compatible fallback


def _set_text(tf, lines, *, size=14, bold=False, italic=False, color=INK,
              font=FONT_BODY, align=PP_ALIGN.LEFT, line_space=1.20):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    for r in list(p.runs):
        r.text = ""
    full = "\n".join(lines)
    p.text = full
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    for para in tf.paragraphs:
        para.line_spacing = line_space


def _add_text_box(slide, left, top, width, height, *, fill=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    box.line.fill.background()
    return box


def _rect_glass(slide, left, top, width, height, *, fill, alpha=40000,
                line=None, shadow=True):
    """A rectangle with frosted-glass alpha + soft shadow."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    from pptx.oxml.ns import qn
    solid = s.fill.fore_color._xFill
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is not None:
        if srgb.find(qn("a:alpha")) is None:
            a = etree.SubElement(srgb, qn("a:alpha"))
            a.set("val", str(alpha))
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    if shadow:
        from pptx.oxml.ns import qn
        sp = s._element
        spPr = sp.find(qn("p:spPr"))
        if spPr is not None:
            existing = spPr.find(qn("a:effectLst"))
            if existing is not None:
                spPr.remove(existing)
            effectLst = etree.SubElement(spPr, qn("a:effectLst"))
            outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
            outerShdw.set("blurRad", "50800")
            outerShdw.set("dist", "25400")
            outerShdw.set("dir", "5400000")
            outerShdw.set("algn", "tl")
            outerShdw.set("rotWithShape", "0")
            color = etree.SubElement(outerShdw, qn("a:srgbClr"))
            color.set("val", "000000")
            alpha_e = etree.SubElement(color, qn("a:alpha"))
            alpha_e.set("val", "22000")
    else:
        s.shadow.inherit = False
    return s


def _page_chrome(slide, point_label: str, page_no: int, total=5):
    """Cobalt Grid chrome: inset hairlines, mono tags and pixel steps."""
    _rect_glass(slide, Inches(0.32), Inches(0.20), Inches(12.68), Emu(14000),
                fill=INK, alpha=100000, line=None, shadow=False)
    _rect_glass(slide, Inches(0.32), Inches(7.42), Inches(12.68), Emu(14000),
                fill=INK, alpha=100000, line=None, shadow=False)
    box = _add_text_box(slide, Inches(0.32), Inches(0.27), Inches(8.2), Inches(0.25))
    _set_text(box.text_frame, [point_label.upper()], size=9.5, bold=True,
              color=INK, font=FONT_MONO)
    box = _add_text_box(slide, Inches(11.60), Inches(7.12), Inches(1.40), Inches(0.22))
    _set_text(box.text_frame, [f"0{page_no} / 0{total}"], size=9.5,
              color=INK, font=FONT_MONO, align=PP_ALIGN.RIGHT)
    hint = _add_text_box(slide, Inches(0.32), Inches(7.12), Inches(3.0), Inches(0.22))
    _set_text(hint.text_frame, ["SPATIAL DOMAIN · CO₂"], size=8.5,
              color=MUTED, font=FONT_MONO)
    # Small stair-stepped pixel-glitch accent.
    for i in range(5):
        _rect_glass(slide, Inches(12.45 + i * 0.085), Inches(0.31 + i * 0.055),
                    Inches(0.065), Inches(0.065), fill=INK,
                    alpha=100000, line=None, shadow=False)


def _add_cobalt_grid(slide):
    """Permanent graph-paper canvas used by every Cobalt Grid slide."""
    step = 0.38
    x = 0.0
    while x <= 13.333:
        _rect_glass(slide, Inches(x), Emu(0), Emu(6500), Inches(7.5),
                    fill=GRID, alpha=100000, line=None, shadow=False)
        x += step
    y = 0.0
    while y <= 7.5:
        _rect_glass(slide, Emu(0), Inches(y), Inches(13.333), Emu(6500),
                    fill=GRID, alpha=100000, line=None, shadow=False)
        y += step


def _add_pic(slide, path, left, top, width, height):
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


# ──────────────────────────────────────────────────────────────────────────────
# Build the deck
# ──────────────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]

ROOT = Path(__file__).resolve().parents[2]
FIGS = ROOT / "1_Exploration of Dataset" / "Figures" / "_Generated_by_Scripts"
P1 = FIGS / "point1"
P2 = FIGS / "point2"
P3 = FIGS / "point3"


# ═══════════════════════════════════════════════════════════════════════════
# Figure page — image on the LEFT (big), frosted-glass card on the RIGHT
# Each slide has ONE BIG TITLE that is the main point of that slide
# ═══════════════════════════════════════════════════════════════════════════
def _figure_page_lr(page_no, point_label, big_title, sub_title, bullets, img_path,
                    img_w=9.10, right_caption=None, numbered_items=None,
                    hide_subtitle_above=False):
    """Layout:
      • slim green header
      • BIG TITLE (the main point) — very large, prominent
      • sub-title (one short line of context) — optionally hidden if
        `hide_subtitle_above=True` (in that case the right card carries it)
      • image on the left (big, dedicated to the figure)
      • frosted-glass content card on the right with:
            1. optional figure caption (the description)
            2. optional numbered items (e.g. small countries in the treemap)
            3. key points
    """
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAGE
    _add_cobalt_grid(s)
    _page_chrome(s, point_label, page_no, total=5)

    # ── BIG TITLE block (above the image, left half) ─────────────────────
    title_left = Inches(0.32)
    title_top = Inches(0.62)
    title_w = Inches(img_w) - Inches(0.10)
    title_h = Inches(0.75)
    box = _add_text_box(s, title_left, title_top, title_w, title_h)
    _set_text(box.text_frame, [big_title], size=21, bold=False, italic=True,
              color=INK, font=FONT_HEAD, line_space=1.04)
    # Sub-title (small one-liner, may wrap to 2 lines) — only when the
    # right card isn't carrying it.
    if sub_title and not hide_subtitle_above:
        box2 = _add_text_box(s, title_left, title_top + title_h,
                             title_w, Inches(0.62))
        _set_text(box2.text_frame, [sub_title], size=11.0, color=INK_SOFT,
                  font=FONT_BODY, line_space=1.20)

    # ── Image card (left) ─────────────────────────────────────────────────
    img_left = Inches(0.30)
    # If the sub-title is hidden (moved to the right card), the image can
    # start higher and become taller — "left side is just a big figure".
    img_top = Inches(1.50) if hide_subtitle_above else Inches(2.05)
    img_h = Inches(5.48) if hide_subtitle_above else Inches(4.93)
    _rect_glass(s, img_left, img_top, Inches(img_w), img_h,
                fill=PAGE, alpha=100000, line=INK, shadow=False)
    with Image.open(img_path) as im:
        iw, ih = im.size
    ratio = iw / ih
    inner_w = Inches(img_w) - Inches(0.10)
    inner_h = img_h - Inches(0.10)
    if ratio > inner_w / inner_h:
        w = inner_w; h = w / ratio
    else:
        h = inner_h; w = h * ratio
    ix = img_left + (Inches(img_w) - w) / 2
    iy = img_top + Inches(0.05) + (inner_h - h) / 2
    _add_pic(s, img_path, ix, iy, w, h)

    # ── Right frosted-glass content card ─────────────────────────────────
    rx = img_left + Inches(img_w) + Inches(0.12)
    rw = SW - rx - Inches(0.30)
    rt = Inches(0.62) if hide_subtitle_above else Inches(0.78)
    # When the right side is the figure's home (caption + numbered items +
    # key points), let the card reach close to the slide bottom — the
    # figure's left card already owns most of the visual area, so the
    # right card is content-driven and can use the full vertical space.
    if hide_subtitle_above:
        rh = Inches(SH / Inches(1) - rt / Inches(1) - 0.20)  # to ~7.10"
    else:
        rh = Inches(6.20)

    _rect_glass(s, rx, rt, rw, rh, fill=PAGE_2, alpha=100000,
                line=INK, shadow=False)
    _rect_glass(s, rx, rt, Inches(0.06), rh,
                fill=ACCENT, alpha=15000, line=None, shadow=False)
    _rect_glass(s, rx + Inches(0.06), rt, rw - Inches(0.06), Inches(0.025),
                fill=INK, line=None, shadow=False)

    # ----- right card content (top → bottom) -----
    inner_x = rx + Inches(0.30)
    inner_w_text = rw - Inches(0.55)
    cur_y = rt + Inches(0.28)

    # (1) Figure description / sub_title (carried over from above when
    #     hide_subtitle_above=True, or supplied explicitly as right_caption)
    cap_text = right_caption if right_caption is not None else (
        sub_title if hide_subtitle_above else None
    )
    if cap_text:
        b_cap = _add_text_box(s, inner_x, cur_y, inner_w_text, Inches(1.55))
        _set_text(b_cap.text_frame, [cap_text], size=11.2, color=INK_SOFT,
                  font=FONT_BODY, line_space=1.30)
        for p in b_cap.text_frame.paragraphs:
            p.space_after = Pt(2)
        cur_y += Inches(1.42)

        # Divider under the caption (thin line)
        _rect_glass(s, inner_x, cur_y + Inches(0.02), inner_w_text, Emu(6000),
                    fill=GRID, alpha=40000, line=None, shadow=False)
        cur_y += Inches(0.18)

    # (2) Numbered items — for the treemap: list of small countries.  The
    #     list is rendered as a 2-column grid so it fits in the narrow card.
    if numbered_items:
        n_hdr = _add_text_box(s, inner_x, cur_y, inner_w_text, Inches(0.28))
        _set_text(n_hdr.text_frame,
                  [f"Numbered countries  ({len(numbered_items)})"],
                  size=11, bold=True, color=ACCENT, font=FONT_HEAD)
        cur_y += Inches(0.30)

        col_w = inner_w_text / 2
        row_h = Inches(0.14)
        rows_per_col = 20  # up to 40 entries fit; the current treemap has 39
        for i, item in enumerate(numbered_items):
            if isinstance(item, tuple):
                num, text = item
            else:
                num, text = i + 1, str(item)
            col = 0 if i < rows_per_col else 1
            row = i if i < rows_per_col else i - rows_per_col
            tx = inner_x + (col_w * col)
            ty = cur_y + (row_h * row)
            tb = _add_text_box(s, tx, ty, col_w, row_h)
            tb.text_frame.word_wrap = False
            p = tb.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run_n = p.add_run(); run_n.text = f"{num:>2}  "
            run_n.font.name = FONT_BODY
            run_n.font.size = Pt(7.5)
            run_n.font.bold = True
            run_n.font.color.rgb = ACCENT
            run_t = p.add_run(); run_t.text = text
            run_t.font.name = FONT_BODY
            run_t.font.size = Pt(7.5)
            run_t.font.color.rgb = INK_SOFT

        n_rows_used = max(min(len(numbered_items), rows_per_col),
                          min(max(0, len(numbered_items) - rows_per_col),
                              rows_per_col))
        cur_y += row_h * n_rows_used + Inches(0.05)

        _rect_glass(s, inner_x, cur_y + Inches(0.02), inner_w_text, Emu(6000),
                    fill=GRID, alpha=40000, line=None, shadow=False)
        cur_y += Inches(0.15)

    # (3) Key points — always present
    kp_top = cur_y
    kp_h = (rt + rh) - kp_top - Inches(0.20)
    b1 = _add_text_box(s, inner_x, kp_top, Inches(1.32), Inches(0.28),
                       fill=INK)
    b1.text_frame.margin_left = Inches(0.08)
    b1.text_frame.margin_right = Inches(0.05)
    b1.text_frame.margin_top = Inches(0.015)
    b1.text_frame.margin_bottom = Inches(0.015)
    _set_text(b1.text_frame, ["KEY POINTS"], size=9.5, bold=True,
              color=PAGE, font=FONT_MONO)
    b2 = _add_text_box(s, inner_x, kp_top + Inches(0.32), inner_w_text,
                       kp_h - Inches(0.32))
    body_lines = [f"•  {h}  {body}" for h, body in bullets]
    _set_text(b2.text_frame, body_lines, size=12.0, color=INK_SOFT,
              font=FONT_BODY, line_space=1.32)
    for p in b2.text_frame.paragraphs:
        p.space_after = Pt(5)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 1 — Intro / story hand-off
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_intro():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(0xEC, 0xEF, 0xF3)
    _page_chrome(s, "Spatial domain  ·  emissions story", 1, total=6)

    # BIG TITLE
    box = _add_text_box(s, Inches(0.32), Inches(0.62), Inches(12.7), Inches(0.75))
    _set_text(box.text_frame,
              ["Major emitters are not a single type of country."],
              size=20, bold=True, color=INK, font=FONT_HEAD, line_space=1.10)
    # Sub-title
    box2 = _add_text_box(s, Inches(0.32), Inches(1.40), Inches(12.7), Inches(0.62))
    _set_text(box2.text_frame,
              ["From \"Who emits the most?\"  to  \"What type of emitter is each country?\"  "
               "— the spatial block proves \"the world is not flat\"."],
              size=11.5, color=INK_SOFT, font=FONT_BODY, line_space=1.20)

    # Three cards: P1 / P2 / P3
    cards = [
        ("P1  ·  Concentration", "Tonnage is highly concentrated",
         ["Top 20 countries = 81% of world CO₂.",
          "Six regions, six colours — Asia is 60% of all emissions."]),
        ("P1 / P2  ·  Scale vs intensity", "Big emitter ≠ dirty emitter",
         ["China leads in tonnage but only 7.9 t / person.",
          "Qatar leads in per-capita at 37+ t / person.",
          "Top-5 tonnage and Top-5 per-capita share only 1 country."]),
        ("P3  ·  Fuel mix", "Fuel mix is a third geography",
         ["190 countries fall into 3 dominant fuels (coal / oil / gas).",
          "Continents pick a fuel; within a continent, neighbours still differ."]),
        ("P3  ·  Trajectory", "Emissions are moving east",
         ["2013 → 2021: Asia +15% (flat per-cap), West down on both axes.",
          "Per-capita: West decoupling; Asia only flat."]),
    ]
    card_xs = [Inches(0.32), Inches(3.50), Inches(6.68), Inches(9.86)]
    card_w = Inches(3.00)
    card_y = Inches(2.20)
    card_h = Inches(4.40)
    for x, (kicker, headline, body) in zip(card_xs, cards):
        # Frosted glass card
        _rect_glass(s, x, card_y, card_w, card_h,
                    fill=RGBColor(0xFF, 0xFF, 0xFF), alpha=82000,
                    line=RGBColor(0xCF, 0xD6, 0xE0))
        # Left accent stripe
        _rect_glass(s, x, card_y, Inches(0.06), card_h,
                    fill=ACCENT, alpha=15000, line=None, shadow=False)
        # Kicker
        bk = _add_text_box(s, x + Inches(0.20), card_y + Inches(0.20),
                           card_w - Inches(0.30), Inches(0.30))
        _set_text(bk.text_frame, [kicker], size=11, bold=True,
                  color=ACCENT, font=FONT_HEAD)
        # Headline
        bh = _add_text_box(s, x + Inches(0.20), card_y + Inches(0.55),
                           card_w - Inches(0.30), Inches(1.10))
        _set_text(bh.text_frame, [headline], size=15, bold=True,
                  color=INK, font=FONT_HEAD, line_space=1.15)
        # Body bullets
        bb = _add_text_box(s, x + Inches(0.20), card_y + Inches(1.70),
                           card_w - Inches(0.30), Inches(2.50))
        body_lines = [f"•  {b}" for b in body]
        _set_text(bb.text_frame, body_lines, size=11.5,
                  color=INK_SOFT, font=FONT_BODY, line_space=1.40)
        for p in bb.text_frame.paragraphs:
            p.space_after = Pt(6)

    # Footer
    fig_text = _add_text_box(s, Inches(0.32), Inches(6.80), Inches(12.7), Inches(0.30))
    _set_text(fig_text.text_frame,
              ["Story line  ·  P1  Who emits the most? → P1  How concentrated?  ·  "
               "P2  Big vs dirty?  ·  P3  What do they burn?  ·  P3  Where is it going?"],
              size=10.5, color=MUTED, font=FONT_BODY)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 2 — P1_03 treemap (region hierarchy)
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_2():
    _figure_page_lr(
        page_no=1,
        point_label="P1  ·  Concentration",
        big_title="A few countries dominate the global total — Asia alone contributes 59.8%.",
        sub_title=("Story step 1 · Who emits the most? Rectangle size = share of world CO₂, "
                   "2021; the hierarchy separates region from country."),
        bullets=[
            ("Concentration", "the Top 20 account for 81.2% of global CO₂; Asia alone contributes 59.8%."),
            ("Reading", "short country names stay in full; long small-country names use the dataset ISO-3 code."),
            ("Bridge", "tonnage is concentrated, but it does not tell us which populations emit most intensively."),
        ],
        img_path=P1 / "P1_03_combo_top20_bars_cumshare_cobalt.png",
        right_caption=("The treemap groups countries by region; rectangle size = share of "
                       "world CO₂, 2021. Legend fill exactly matches country cells; the darker "
                       "outline marks the region boundary. No separate country index is needed."),
        numbered_items=None,
        hide_subtitle_above=True,
    )


# ═══════════════════════════════════════════════════════════════════════════
# Slide 3 — P1 3D map (legend below)
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_3():
    _figure_page_lr(
        page_no=2,
        point_label="P1  ·  Tonnage vs intensity",
        big_title="Big emitter ≠ high per-capita emitter — Singapore sits between the two extremes.",
        sub_title="Story step 2 · Land colour = log tonnage; 3-D bar height/colour = per-capita CO₂. Major emitters, Qatar and Singapore are labelled without overlap.",
        bullets=[
            ("Two extremes", "China leads total emissions, while tiny Qatar has the tallest per-capita bar."),
            ("Singapore", "50.1 Mt and 9.03 t / person — included as local context with the same styling as other countries."),
            ("Bridge", "the map motivates a direct country-by-country comparison of scale and intensity."),
        ],
        img_path=P1 / "P1_01_map_bars_total_co2_2021_cobalt.png",
    )


# ═══════════════════════════════════════════════════════════════════════════
# Slide 4 — P2 bubble (all countries, top-5 labels, 2-D gradient colour)
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_4():
    _figure_page_lr(
        page_no=3,
        point_label="P2  ·  Scale vs intensity",
        big_title="High total ≠ high per-capita — the two Top-5 lists do not overlap.",
        sub_title=("Story step 3 · Every country is a circle: x = log₁₀ total CO₂, "
                   "y = per-capita CO₂, size = population, and colour blends both axes."),
        bullets=[
            ("Different leaders", "the total-emissions and per-capita Top-5 lists share no country (0 / 5)."),
            ("Spearman ρ = 0.36", "weak positive correlation — being big slightly raises per-capita, but the scatter is wide."),
            ("Take-away", "‘big emitter’ ≠ ‘high per-capita emitter’ — the two groups occupy different zones."),
        ],
        img_path=P2 / "P2_02_combo_bubble_scale_vs_intensity_cobalt.png",
    )


# ═══════════════════════════════════════════════════════════════════════════
# Slide 5 — P3 fuel map (light bg, proper pie proportions, larger map)
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_5():
    _figure_page_lr(
        page_no=4,
        point_label="P3  ·  Fuel geography",
        big_title="Fuel mix forms spatial regimes — a natural bridge from maps to country clustering.",
        sub_title=("Story step 4 · Country fill = dominant fuel (Coal / Oil / Gas); "
                   "top emitters carry real 2021 fuel-mix pies. The compact legend gives the map more room."),
        bullets=[
            ("Spatial regimes", "coal clusters in Asia; oil spans Africa, the Middle East and the Americas; gas is prominent across Russia and Central Asia."),
            ("Real pies", "show the truth behind the colour: e.g. China is 78% coal, but Indonesia and S. Arabia look oil-red even when their mix is mixed."),
            ("Next step", "coal/oil/gas shares provide interpretable features for the later K-means grouping of emission patterns."),
        ],
        img_path=P3 / "P3_01_map_dominant_fuel_2021_cobalt.png",
    )


# ═══════════════════════════════════════════════════════════════════════════
# Slide 6 — P3 trajectory (4 chains × 5 time points = 20 circles)
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_6():
    _figure_page_lr(
        page_no=5,
        point_label="P3  ·  Time–space trajectory 2013 → 2021",
        big_title="Total-emissions geography moves east, while per-capita trends tell a different story.",
        sub_title="Story step 5 · Five two-year observations per track. Circle area grows with time and track scale; semi-transparent arrows show 2013→2021 direction only, not distance.",
        bullets=[
            ("Method", "for each year, countries are placed at their map representative points and averaged using national CO₂ emissions as weights."),
            ("Direction", "the global centre shifts east-south-east by about 585 km; Africa moves north-west, South America south-west, while North America is nearly stationary."),
            ("Conclusion", "the movements are modest, not continental jumps. The main spatial story is a gradual eastward rebalancing as Asia rises and Europe declines."),
        ],
        img_path=P3 / "P3_06_trajectory_pseudo3d_cobalt.png",
    )


# Build
build_slide_2()
build_slide_3()
build_slide_4()
build_slide_5()
build_slide_6()

out = (
    ROOT
    / "code"
    / "_cleanup_archive_20260827"
    / "ppt_intermediates"
    / "空间域_PPT_5页_CobaltGrid_修订版.pptx"
)
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"saved {out}")
