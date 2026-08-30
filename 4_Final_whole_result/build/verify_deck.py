# -*- coding: utf-8 -*-
"""Quick QA on the built deck: fonts, minimum sizes, animation and page count."""
import re
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

DECK = Path(r"C:\Users\user\Desktop\Overall_Data Visual Assignment"
            r"\4_Final_whole_result\CA6002_Group30_Final_Presentation.pptx")

z = zipfile.ZipFile(DECK)
slides = sorted(x for x in z.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", x))
anim = sum(1 for x in slides if b"<p:timing" in z.read(x))
trans = sum(1 for x in slides if b"<p:transition" in z.read(x))
faces = {}
for x in slides:
    for m in re.finditer(r'typeface="([^"]+)"', z.read(x).decode("utf8")):
        faces[m.group(1)] = faces.get(m.group(1), 0) + 1
print("slides:", len(slides), "| with animation:", anim, "| with transition:", trans)
print("typefaces used:", faces)

prs = Presentation(str(DECK))
small = []
for i, s in enumerate(prs.slides, 1):
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        tight = (sh.width or 0) < Inches(1.75) or (sh.height or 0) < Inches(0.3)
        for p in sh.text_frame.paragraphs:
            sizes = [r.font.size.pt for r in p.runs if r.font.size is not None]
            if p.font.size is not None:
                sizes.append(p.font.size.pt)
            for v in sizes:
                if v < (10.0 if tight else 11.0):
                    small.append((i, v, p.text[:40]))
print("paragraphs below the size floor:", len(small))
for row in small[:20]:
    print("   ", row)
print("banners (stage : conclusion-style title  |  counter):")
FRONT_BACK = {1, 2}          # cover and contents
for i, s in enumerate(prs.slides, 1):
    band = [sh.text_frame.text.strip() for sh in s.shapes
            if sh.has_text_frame and sh.top is not None and sh.top < Inches(0.9)
            and sh.text_frame.text.strip()]
    if not band:
        print(f"  {i:>2}: (no banner text)")
        continue
    block = max(band, key=len)
    counter = next((t for t in band if t != block), "-")
    lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
    kicker = lines[0] if len(lines) > 1 else ""
    title = lines[-1]
    body = i not in FRONT_BACK and i != len(prs.slides._sldIdLst)
    flag = "   << context line missing" if body and " \u00b7 " not in kicker else ""
    print(f"  {i:>2}: {kicker[:74]:<74} | {counter}")
    print(f"      {title[:74]}{flag}")
