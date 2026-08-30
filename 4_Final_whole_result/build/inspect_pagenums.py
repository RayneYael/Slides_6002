"""List every shape in the bottom-right corner that looks like a page number."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

DECK = Path(__file__).resolve().parents[1] / "CA6002_Group30_Final_Presentation.pptx"
prs = Presentation(str(DECK))

for i, s in enumerate(prs.slides, 1):
    for sh in s.shapes:
        if not sh.has_text_frame or sh.left is None or sh.top is None:
            continue
        t = sh.text_frame.text.strip()
        if sh.left > Inches(11.0) and sh.top > Inches(6.3) and len(t) <= 4:
            ph = sh.is_placeholder and sh.placeholder_format.type
            print(f"slide {i:>2}  x={sh.left / 914400:.2f} y={sh.top / 914400:.2f} "
                  f"w={sh.width / 914400:.2f} text={t!r} placeholder={ph} name={sh.name!r}")
