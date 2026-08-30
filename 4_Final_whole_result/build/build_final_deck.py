# -*- coding: utf-8 -*-
"""Build the complete CA6002 Group 30 final deck.

Structure (cover / contents / acknowledgements are not counted in the 20 body
slides, per the assignment template):

    cover · contents
    1      Introduction
    2-8    Exploration of Dataset          (7 slides imported from the source deck)
    9-12   Design of AI Algorithm
    13-15  Model Evaluation and Performance
    16-19  Visual Storytelling             (all four are figure-led)
    20     Conclusions
    acknowledgements

Design system
  * Times New Roman everywhere (theme fonts + every run), minimum body size 11 pt.
  * Every slide opens with a left-aligned gradient banner carrying the assignment
    stage name, a colon, and the slide's conclusion-style title.
  * Low-saturation palette shared with make_figures.py: muted steel blue for the
    plateauing pathway, muted amber for the rapid-growth pathway, deep maroon as
    the only accent colour.
  * The NTU template background (red rule + logo) sits on every slide.
"""
import os
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

ROOT = Path(r"C:\Users\user\Desktop\Overall_Data Visual Assignment")
# build input, deliberately kept inside build/ so the delivery folder holds only
# the finished deck (see build/source/README.txt)
SRC = ROOT / r"4_Final_whole_result/build/source/exploration_slides_source.pptx"
# DECK_OUT lets a build go to a scratch file while the deliverable is open in
# PowerPoint, which otherwise locks it against writing
OUT = Path(os.environ.get("DECK_OUT")) if os.environ.get("DECK_OUT") else \
    ROOT / r"4_Final_whole_result/CA6002_Group30_Final_Presentation.pptx"
FIG = ROOT / r"4_Final_whole_result/build/figures"
ASSETS = ROOT / r"4_Final_whole_result/build/assets"
BG = ASSETS / "image1.png"          # NTU template background: red rule + logo

FONT = "Times New Roman"          # body copy: the academic default
FONT_TITLE = "Georgia"            # banners, display numbers: serif but distinct
MIN_PT = 13.0
BUMP = 2.0                        # global type-size lift, applied in tb()/banner()

ACCENT = RGBColor(0x2F, 0x52, 0x48)       # deep forest green, the chrome colour
ACCENT_L = RGBColor(0x76, 0x96, 0x82)     # banner gradient end
INK = RGBColor(0x25, 0x2B, 0x30)
BLUE = RGBColor(0x4F, 0x6E, 0x96)         # plateauing / slowdown pathway
ORANGE = RGBColor(0xBE, 0x7F, 0x4E)       # rapid-growth pathway
LBLUE = RGBColor(0xEC, 0xEF, 0xF4)
LORANGE = RGBColor(0xF7, 0xF1, 0xEA)
GREY = RGBColor(0xFA, 0xFB, 0xFA)         # frosted card fill (used with alpha)
MIDGREY = RGBColor(0x6E, 0x76, 0x72)
LINEGREY = RGBColor(0xD5, 0xDE, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BANNER_SUB = RGBColor(0xD3, 0xE3, 0xD6)   # banner sub-label
WASH_A = RGBColor(0xFF, 0xFF, 0xFF)       # slide background: white to pale green
WASH_B = RGBColor(0xEC, 0xF3, 0xEA)
NTU_RED = RGBColor(0xC8, 0x10, 0x2E)      # template footer rule only

GLASS_ALPHA = 74000               # 74% opaque frosted cards
CARD_FILLS = {GREY, WHITE, LBLUE, LORANGE}

GENAI = ("Generative AI acknowledgement: ChatGPT (OpenAI) assisted with Python "
         "code for figure generation and slide assembly; all data processing, model "
         "results, and interpretation are the group's own work.")

STAGE_EXPLORE = "EXPLORATION OF DATASET"
STAGE_DESIGN = "DESIGN OF AI ALGORITHM"
STAGE_EVAL = "MODEL EVALUATION AND PERFORMANCE"
STAGE_STORY = "VISUAL STORYTELLING"

ANIM = {}
_counter = [0]

# Two playback modes for the whole deck:
#   AUTO_PLAY = True   every group starts by itself once the previous one ends,
#                      so a slide plays through without a single click
#   AUTO_PLAY = False  body-slide groups wait for a click (presenter-paced)
# Slide-to-slide advance stays manual either way, so the speaker still decides
# when to leave a page.
AUTO_PLAY = True
AUTO_GAP = 400          # ms of breathing room before an auto group starts
AUTO_STAGGER = 130      # ms between shapes inside one auto group


def tag(slide, *shapes, auto=False, delay=0, stagger=0, filt="fade", dur=450):
    """Register an animation group.

    auto=False  the group waits for a click, unless AUTO_PLAY overrides it
    auto=True   the group plays on its own, `delay` ms after the previous group
                ends, and each shape inside it is offset by a further `stagger` ms
    """
    if AUTO_PLAY and not auto:
        auto, delay, stagger = True, AUTO_GAP, AUTO_STAGGER
    names = []
    for sh in shapes:
        _counter[0] += 1
        sh.name = f"anim_{_counter[0]}"
        names.append(sh.name)
    ANIM.setdefault(id(slide._element), []).append(
        {"names": names, "auto": auto, "delay": delay, "stagger": stagger,
         "filt": filt, "dur": dur})


# ----------------------------------------------------------------- motion
def build_timing_xml(groups):
    """PowerPoint-standard mainSeq.

    A group is either a click group (clickEffect, waits for the presenter) or an
    automatic group that starts with the slide and staggers its own shapes.
    """
    P = 'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    nid = [3]

    def nidn():
        nid[0] += 1
        return nid[0]

    def effect_par(spid, node_type, delay, filt, dur):
        i1 = nidn()
        return (
            f'<p:par><p:cTn id="{i1}" presetID="10" presetClass="entr" presetSubtype="0" '
            f'fill="hold" grpId="0" nodeType="{node_type}">'
            f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst><p:childTnLst>'
            f'<p:set><p:cBhvr><p:cTn id="{nidn()}" dur="1" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
            f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
            f'<p:animEffect transition="in" filter="{filt}"><p:cBhvr>'
            f'<p:cTn id="{nidn()}" dur="{dur}"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>'
            f'</p:childTnLst></p:cTn></p:par>')

    clicks = ""
    for grp in groups:
        auto = grp["auto"]
        start = "0" if auto else "indefinite"
        inner = "".join(
            effect_par(spid,
                       ("afterEffect" if j == 0 else "withEffect") if auto
                       else ("clickEffect" if j == 0 else "withEffect"),
                       (grp["delay"] + j * grp["stagger"]) if auto else 0,
                       grp["filt"], grp["dur"])
            for j, spid in enumerate(grp["spids"]))
        clicks += (f'<p:par><p:cTn id="{nidn()}" fill="hold">'
                   f'<p:stCondLst><p:cond delay="{start}"/></p:stCondLst>'
                   f'<p:childTnLst><p:par><p:cTn id="{nidn()}" fill="hold">'
                   f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
                   f'<p:childTnLst>{inner}</p:childTnLst></p:cTn></p:par>'
                   f'</p:childTnLst></p:cTn></p:par>')
    return (
        f'<p:timing {P}><p:tnLst><p:par>'
        f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>{clicks}'
        f'</p:childTnLst></p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>')


TRANSITION_XML = ('<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                  'spd="med"><p:fade/></p:transition>')


def inject_motion():
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    for slide in prs.slides:
        el = slide._element
        for old in el.findall(f"{{{ns}}}transition"):
            el.remove(old)
        el.append(parse_xml(TRANSITION_XML))
    for slide in prs.slides:
        groups = ANIM.get(id(slide._element))
        if not groups:
            continue
        name2id = {sh.name: sh.shape_id for sh in slide.shapes}
        resolved = []
        for grp in groups:
            spids = [name2id[n] for n in grp["names"] if n in name2id]
            if spids:
                resolved.append(dict(grp, spids=spids))
        if resolved:
            slide._element.append(parse_xml(build_timing_xml(resolved)))


prs = Presentation(str(SRC))
LAYOUT = prs.slide_layouts[0]
SW_IN = prs.slide_width / 914400


# ----------------------------------------------------------------- helpers
def tb(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
       fill=None, line_color=None, wrap=True, line_spacing=None):
    """paras: list of dict(runs=[(text,size,bold,color)], align=, space_before=)."""
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sh.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", align)
        if line_spacing:
            p.line_spacing = line_spacing
        if para.get("space_before"):
            p.space_before = Pt(para["space_before"])
        for (t, size, bold, color) in para["runs"]:
            r = p.add_run()
            r.text = t
            r.font.name = FONT_TITLE if size >= 24 else FONT
            r.font.size = Pt(size + (BUMP if size < 24 else 0))
            r.font.bold = bold
            r.font.color.rgb = color
    if fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color is not None:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def wordart(slide, x, y, w, h, text, subscript, size):
    """Display lettering for the cover: gradient-filled, outlined, softly shadowed.

    python-pptx exposes none of these run effects, so the run properties are
    written directly; the subscript run carries the chemical "2" of CO2.
    """
    A = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.background()
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for body, base, pt in ((text, None, size), (subscript, -25000, size * 0.62)):
        r = p.add_run()
        r.text = body
        rPr = r._r.get_or_add_rPr()
        for child in list(rPr):
            rPr.remove(child)
        rPr.set("sz", str(int(pt * 100)))
        rPr.set("b", "1")
        if base is not None:
            rPr.set("baseline", str(base))
        rPr.append(parse_xml(
            f'<a:ln {A} w="12700"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
            '</a:ln>'))
        rPr.append(parse_xml(
            f'<a:gradFill {A}><a:gsLst>'
            '<a:gs pos="0"><a:srgbClr val="769682"/></a:gs>'
            '<a:gs pos="100000"><a:srgbClr val="2F5248"/></a:gs>'
            '</a:gsLst><a:lin ang="5400000" scaled="0"/></a:gradFill>'))
        rPr.append(parse_xml(
            f'<a:effectLst {A}><a:outerShdw blurRad="50800" dist="25400" dir="2700000" '
            'algn="tl" rotWithShape="0"><a:srgbClr val="1F2A24"><a:alpha val="20000"/>'
            '</a:srgbClr></a:outerShdw></a:effectLst>'))
        rPr.append(parse_xml(f'<a:latin {A} typeface="{FONT_TITLE}"/>'))
    return sh


def set_alpha(shape, alpha=GLASS_ALPHA):
    """python-pptx cannot express fill transparency, so patch the srgbClr node."""
    clr = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
    if clr is None:
        return
    for old in clr.findall(qn("a:alpha")):
        clr.remove(old)
    clr.append(parse_xml(
        '<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        f'val="{alpha}"/>'))


def soft_shadow(shape, blur=44000, dist=14000, alpha=13000):
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    spPr.append(parse_xml(
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{blur}" dist="{dist}" dir="5400000" rotWithShape="0">'
        f'<a:srgbClr val="2F5248"><a:alpha val="{alpha}"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'))


def rect(slide, x, y, w, h, fill, line_color=None, shape=MSO_SHAPE.RECTANGLE,
         glass=None):
    """Panels. Light fills are rendered as frosted glass: translucent, hairline,
    soft drop shadow; solid accent fills (rules, chips) stay flat."""
    if glass is None:
        glass = fill in CARD_FILLS and w > 0.5 and h > 0.25
    if glass and shape is MSO_SHAPE.RECTANGLE:
        shape = MSO_SHAPE.ROUNDED_RECTANGLE
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if glass:
        set_alpha(sh)
        soft_shadow(sh)
        if shape is MSO_SHAPE.ROUNDED_RECTANGLE:
            sh.adjustments[0] = 0.055
        if line_color is None:
            line_color = LINEGREY
    if line_color is not None:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    if not glass:
        sh.shadow.inherit = False
    return sh


def gradient_rect(slide, x, y, w, h, c1, c2, angle=0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.gradient()
    sh.fill.gradient_angle = angle
    stops = sh.fill.gradient_stops
    stops[0].color.rgb = c1
    stops[0].position = 0.0
    stops[1].color.rgb = c2
    stops[1].position = 1.0
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def send_to_back(slide, shape):
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def background(slide, hills=False):
    """White-to-pale-green wash, then the NTU footer redrawn on top of it.

    The template ships its footer as one opaque white picture, which would hide
    the wash, so the rule and the logo are placed as separate elements instead.
    """
    wash = gradient_rect(slide, 0, 0, SW_IN, prs.slide_height / 914400,
                         WASH_A, WASH_B, angle=45)
    send_to_back(slide, wash)
    if hills:
        pic = slide.shapes.add_picture(str(FIG / "art_hills.png"), Inches(0),
                                       Inches(5.55), Inches(SW_IN), Inches(1.4))
        spTree = slide.shapes._spTree
        spTree.remove(pic._element)
        spTree.insert(3, pic._element)
    rule = rect(slide, 0.0, 6.86, SW_IN, 0.028, NTU_RED, glass=False)
    slide.shapes.add_picture(str(ASSETS / "ntu_logo.png"), Inches(0.32),
                             Inches(6.96), Inches(1.24), Inches(0.5))
    return rule


def page_number(slide, n):
    tb(slide, 12.25, 7.10, 0.82, 0.32,
       [{"runs": [(str(n), 12, False, MIDGREY)], "align": PP_ALIGN.RIGHT}],
       align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def banner(slide, stage, title, counter=None, topic=None):
    """Left-aligned gradient banner, read top-down as context then finding.

    Line 1  STAGE · what this page investigates
    Line 2  the conclusion-style title

    The topic line matters: a marker who does not carry our background should not
    meet a bare four-word verdict with no idea what is being investigated.
    """
    bar = gradient_rect(slide, 0, 0, SW_IN, 0.96, ACCENT, ACCENT_L)
    rect(slide, 0, 0.96, SW_IN, 0.022, ACCENT_L, glass=False)
    box = slide.shapes.add_textbox(Inches(0.42), Inches(0.03), Inches(10.9), Inches(0.90))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    kicker = " · ".join(x for x in (stage, topic) if x)
    if kicker:
        r = p.add_run()
        r.text = kicker
        r.font.name = FONT_TITLE
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = BANNER_SUB
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
    r2 = p.add_run()
    r2.text = title
    r2.font.name = FONT_TITLE
    # keep the conclusion line to a single line whatever its length
    r2.font.size = Pt(25 if len(title) <= 46 else
                      (23 if len(title) <= 58 else (21 if len(title) <= 70 else 19)))
    r2.font.bold = True
    r2.font.color.rgb = WHITE
    box.fill.background()
    box.line.fill.background()
    box.shadow.inherit = False
    if counter:
        tb(slide, 11.45, 0.30, 1.55, 0.36,
           [{"runs": [(counter, 12, True, BANNER_SUB)], "align": PP_ALIGN.RIGHT}],
           align=PP_ALIGN.RIGHT)
    return bar


def responsible(slide, names, refs=None):
    # sits low on the footer line, vertically centred against the NTU logo, and
    # always in the same grey small-caps style on every slide
    tb(slide, 1.72, 7.10, 7.0, 0.32,
       [{"runs": [(f"Responsible: {names.upper()}", 11, True, MIDGREY)]}],
       anchor=MSO_ANCHOR.MIDDLE)
    if refs:
        ref_mark(slide, refs)


def ref_mark(slide, refs):
    """Small reference marker on the footer line, keyed to the list on the closing slide."""
    return tb(slide, 8.55, 7.10, 3.0, 0.32,
              [{"runs": [(refs, 9, False, MIDGREY)], "align": PP_ALIGN.RIGHT}],
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def picture_fit(slide, path, x, y, w, h):
    """Insert a picture centred in the (x, y, w, h) box, preserving aspect ratio."""
    im = Image.open(path)
    ar = im.width / im.height
    if ar >= w / h:
        nw, nh = w, w / ar
    else:
        nh, nw = h, h * ar
    return slide.shapes.add_picture(str(path), Inches(x + (w - nw) / 2),
                                   Inches(y + (h - nh) / 2), Inches(nw), Inches(nh))


def caption(slide, x, y, w, text, size=11.5):
    return tb(slide, x, y, w, 0.36, [{"runs": [(text, size, False, MIDGREY)]}])


def label(slide, x, y, w, text):
    return tb(slide, x, y, w, 0.32, [{"runs": [(text, 11.5, True, ACCENT)]}])


def clear_placeholders(slide):
    """Drop the layout placeholders: every slide is composed from explicit shapes."""
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def new_slide(stage, title, counter, resp, number, topic=None, refs=None):
    s = prs.slides.add_slide(LAYOUT)
    background(s)
    clear_placeholders(s)
    banner(s, stage, title, counter, topic)
    responsible(s, resp, refs)
    page_number(s, number)
    return s


# =================================================================
# COVER  (not counted in the 20 body slides)
# =================================================================
s = prs.slides.add_slide(LAYOUT)
background(s, hills=True)
clear_placeholders(s)
gradient_rect(s, 0, 0, SW_IN, 0.34, ACCENT, ACCENT_L)
rect(s, 0.55, 1.05, 0.075, 3.6, ACCENT, glass=False)

# --- title block -------------------------------------------------------------
c_eyebrow = tb(s, 0.85, 1.02, 6.6, 0.36,
               [{"runs": [("CA6002  ·  AI UX & DATA VISUALISATION",
                           12.5, True, ACCENT)]}])
c_title = tb(s, 0.85, 1.38, 6.6, 1.15,
             [{"runs": [("Investigating Country", 32, True, INK)]},
              {"runs": [("Carbon Emission Patterns", 32, True, INK)]}],
             line_spacing=0.95)
c_lead = tb(s, 0.85, 2.52, 6.6, 0.5,
            [{"runs": [("- one rising total, two transition pathways", 19, True, ACCENT)]}])
c_rule = rect(s, 0.85, 3.18, 1.5, 0.045, ACCENT, glass=False)
c_sub = tb(s, 0.85, 3.34, 6.3, 1.55,
           [{"runs": [("Fossil CO2 of 193 countries, 1992-2021: does one shared trajectory "
                       "explain the rise?", 16.5, False, INK)]},
            {"runs": [("The covered total rose from ", 16.5, False, INK),
                      ("21.97 Gt", 16.5, True, ACCENT),
                      (" to ", 16.5, False, INK),
                      ("35.98 Gt", 16.5, True, ACCENT),
                      ("  (+63.8%)", 16.5, True, ACCENT)], "space_before": 7},
            {"runs": [("K-means used as a screening tool for climate-monitoring priorities",
                       14, False, MIDGREY)], "space_before": 7}])

# --- carbon-theme illustration: coal plant on the left, clear sky on the right
c_sky = s.shapes.add_picture(str(FIG / "art_sky.png"), Inches(8.45), Inches(0.55),
                             Inches(4.4), Inches(2.2))
c_plant = s.shapes.add_picture(str(FIG / "art_plant.png"), Inches(1.05), Inches(5.20),
                               Inches(2.70), Inches(1.66))
# plumes are anchored on the three chimney mouths of art_plant
c_smoke = [s.shapes.add_picture(str(FIG / f"art_smoke{i}.png"), Inches(x), Inches(y),
                                Inches(w), Inches(w * 2.3 / 1.5))
           for i, (x, y, w) in enumerate([(2.25, 4.71, 0.52), (2.565, 5.09, 0.45),
                                          (2.83, 5.39, 0.40)], start=1)]

# --- evidence and authorship -------------------------------------------------
c_spark = wordart(s, 7.55, 2.10, 5.0, 1.85, "CO", "2", 132)
c_spark_cap = tb(s, 7.30, 4.14, 5.45, 0.4,
                 [{"runs": [("FOSSIL CARBON DIOXIDE  ·  193 COUNTRIES",
                             12.5, True, MIDGREY)]}], align=PP_ALIGN.CENTER)
c_group_card = rect(s, 7.30, 4.78, 5.45, 1.55, GREY)
c_group = tb(s, 7.52, 4.92, 5.05, 1.3,
             [{"runs": [("Group 30", 19, True, ACCENT)]},
              {"runs": [("RAN CHANGMING  ·  REN SIYU  ·  REYIZA RESIBIEKE  ·  SHEN RUITING",
                         13, True, INK)], "space_before": 5},
              {"runs": [("Nanyang Technological University  ·  Group Assignment",
                         12.5, False, MIDGREY)], "space_before": 4}])

# cover animation plays by itself: title, then sky, plant, rising smoke, evidence
# Auto delays are gaps after the previous group ends, so they stay short and the
# whole cover builds in about five seconds.
tag(s, c_eyebrow, c_title, c_lead, c_rule, c_sub,
    auto=True, delay=200, stagger=240, dur=600)
tag(s, c_sky, auto=True, delay=150, dur=700)
tag(s, c_plant, auto=True, delay=150, filt="wipe(up)", dur=700)
tag(s, *c_smoke, auto=True, delay=150, stagger=300, filt="wipe(up)", dur=900)
tag(s, c_spark, c_spark_cap, auto=True, delay=150, stagger=180, dur=600)
tag(s, c_group_card, c_group, auto=True, delay=150, stagger=180, dur=600)
notes(s, "Title slide. Text-only cover by design: no data graphic here, only the framing question and "
         "the one number the whole deck investigates. Group 30. " + GENAI)

# =================================================================
# CONTENTS  (not counted)
# =================================================================
s = prs.slides.add_slide(LAYOUT)
background(s)
clear_placeholders(s)
banner(s, None, "Contents", topic="How this carbon emission investigation is organised")
tb(s, 0.85, 1.06, 11.6, 0.34,
   [{"runs": [("Six blocks, one question: ", 12.5, True, ACCENT),
              ("covered fossil CO2 kept rising for 30 years - one shared trajectory, or "
               "several transition pathways?", 12.5, False, INK)]}])
rows = [
    ("1", "Introduction", "Slide 3", "The puzzle: 30 years of growth, one path or many?", ACCENT),
    ("2", "Exploration of Dataset", "Slides 4-10",
     "Sources, conditioning, and where countries differ", BLUE),
    ("3", "Design of AI Algorithm", "Slides 11-14",
     "From trajectories to 11 features, PCA and K-means", ORANGE),
    ("4", "Model Evaluation and Performance", "Slides 15-17",
     "Validity, stability and membership uncertainty", BLUE),
    ("5", "Visual Storytelling", "Slides 18-21",
     "Two pathways, two monitoring priorities", ORANGE),
    ("6", "Conclusions", "Slide 22", "Contributions, novelty and limitations", ACCENT),
]
y = 1.48
for num, name, pages, desc, accent in rows:
    rect(s, 0.85, y, 11.6, 0.76, GREY if int(num) % 2 else WHITE)
    rect(s, 0.85, y, 0.07, 0.76, accent)
    tb(s, 1.08, y + 0.10, 0.6, 0.55, [{"runs": [(num, 19, True, accent)]}])
    tb(s, 1.75, y + 0.05, 6.2, 0.66,
       [{"runs": [(name, 16, True, INK)]},
        {"runs": [(desc, 11.5, False, MIDGREY)]}])
    tb(s, 10.3, y + 0.16, 2.0, 0.45,
       [{"runs": [(pages, 12.5, True, accent)], "align": PP_ALIGN.RIGHT}],
       align=PP_ALIGN.RIGHT)
    y += 0.84
tb(s, 0.85, 6.55, 11.6, 0.34,
   [{"runs": [("Each block below names the stage it reports, and every slide banner "
               "repeats that stage name.", 11, False, MIDGREY)]}])
notes(s, "Contents slide. Section names match the banner on every later slide, so the "
         "audience can map any page back to this overview. " + GENAI)

# =================================================================
# BODY 1 — INTRODUCTION
# =================================================================
s = new_slide("INTRODUCTION", "Emissions kept rising - but not along one shared path",
              None, "SHEN RUITING", 3,
              topic="Why we investigate country carbon emission patterns")
label(s, 0.5, 1.10, 3.0, "THE PUZZLE")
c1 = rect(s, 0.5, 1.42, 3.6, 1.75, LBLUE)
t1 = tb(s, 0.7, 1.56, 3.2, 1.5,
        [{"runs": [("+63.8%", 28, True, BLUE)]},
         {"runs": [("Sample-wide fossil CO2, 1992 to 2021", 12, True, INK)], "space_before": 4},
         {"runs": [("21.97 to 35.98 Gt across the 193 countries covered by our panel",
                    11, False, INK)], "space_before": 3}])
c2 = rect(s, 0.5, 3.35, 3.6, 1.75, LORANGE)
t2 = tb(s, 0.7, 3.49, 3.2, 1.5,
        [{"runs": [("30 years", 28, True, ORANGE)]},
         {"runs": [("of sustained global climate attention", 12, True, INK)], "space_before": 4},
         {"runs": [("yet the covered total kept growing - why, and driven by whom?",
                    11, False, INK)], "space_before": 3}])
tag(s, c1, t1)
tag(s, c2, t2)
label(s, 4.6, 1.10, 6.0, "OUR QUESTION CHAIN")
questions = [
    "Is the rise one common trajectory - or different transition pathways?",
    "Do countries differ only in size, or in speed, timing and fuel mix?",
    "Can AI identify these pathways stably from trajectory features?",
    "Where should limited climate-monitoring attention go?",
]
y = 1.42
for i, q in enumerate(questions):
    accent = BLUE if i < 3 else ACCENT
    r1 = rect(s, 4.6, y, 7.85, 0.86, WHITE, line_color=accent)
    r2 = tb(s, 4.82, y + 0.06, 0.62, 0.74, [{"runs": [(f"Q{i+1}", 15, True, accent)]}],
            anchor=MSO_ANCHOR.MIDDLE)
    r3 = tb(s, 5.45, y + 0.06, 6.85, 0.74,
            [{"runs": [(q, 13.5, i == 3, INK)]}], anchor=MSO_ANCHOR.MIDDLE)
    tag(s, r1, r2, r3)
    y += 0.98
b1 = rect(s, 0.5, 5.42, 11.95, 1.25, GREY)
b2 = tb(s, 0.72, 5.52, 11.5, 1.1,
        [{"runs": [("Approach", 12, True, ACCENT)]},
         {"runs": [("explore differences  \u2192  engineer trajectory features  \u2192  "
                    "K-means (K = 2)  \u2192  evaluate separation and stability  \u2192  "
                    "pathway-specific monitoring priorities", 13, False, INK)],
          "space_before": 4},
         {"runs": [("Each stage below is reported in the order it was carried out.",
                    11, False, MIDGREY)], "space_before": 4}])
tag(s, b1, b2)
notes(s, "Motivation: despite three decades of climate attention, the sample-wide fossil CO2 "
         "total still grew 63.8% between 1992 and 2021 (21.97 to 35.98 Gt). This slide poses the "
         "investigative question chain that structures the whole deck; the maroon-framed Q4 is the "
         "managerial question the Visual Storytelling section answers. Design rationale: two "
         "high-salience stat cards establish the contradiction before any methodology, and the "
         "numbered chain previews the four analytical stages. " + GENAI)

# =================================================================
# BODY 9-12 — DESIGN OF AI ALGORITHM
# =================================================================
s = new_slide(STAGE_DESIGN, "From 213 territories to 193 comparable country profiles",
              "1 / 4", "RAN CHANGMING", 11,
              topic="Defining the clustering unit and the modelled cohort")
label(s, 0.5, 1.08, 6.0, "COHORT FUNNEL")
funnel = [
    ("6,378", "country-year rows", "OWID backbone + GCB fuel fills, 1992-2021", LBLUE, BLUE),
    ("213", "countries & territories", "valid ISO codes in the panel", LBLUE, BLUE),
    ("-16", "micro states excluded", "populations too small for stable trajectories", LORANGE, ORANGE),
    ("-4", "incomplete histories excluded", "fewer than 30 consecutive years 1992-2021", LORANGE, ORANGE),
    ("193", "modelled countries / entities", "complete history + reliable fuel structure", GREY, ACCENT),
]
y = 1.42
for i, (big, lab_, sub, fill, accent) in enumerate(funnel):
    indent = 0.32 * i
    w = 6.4 - 0.45 * i
    r1 = rect(s, 0.5 + indent, y, w, 0.94, fill)
    r2 = tb(s, 0.66 + indent, y + 0.09, 1.4, 0.75, [{"runs": [(big, 20, True, accent)]}])
    r3 = tb(s, 2.06 + indent, y + 0.07, w - 1.7, 0.82,
            [{"runs": [(lab_, 12.5, True, INK)]},
             {"runs": [(sub, 11, False, MIDGREY)]}])
    tag(s, r1, r2, r3)
    y += 1.06
label(s, 7.35, 1.08, 5.4, "WHY THIS COHORT")
r1 = rect(s, 7.35, 1.42, 5.1, 2.05, WHITE, line_color=BLUE)
r2 = tb(s, 7.55, 1.56, 4.7, 1.8,
        [{"runs": [("Trajectory features need complete histories. ", 12, True, INK),
                   ("Slopes, peaks and post-peak change are only comparable when every country "
                    "covers the same 30 years.", 12, False, INK)]},
         {"runs": [("Fuel features need reliable shares. ", 12, True, INK),
                   ("Entities with unreliable fuel structure would distort the fuel-shift block.",
                    12, False, INK)], "space_before": 8}])
tag(s, r1, r2)
r3 = rect(s, 7.35, 3.72, 5.1, 1.6, GREY)
r4 = tb(s, 7.55, 3.86, 4.7, 1.35,
        [{"runs": [("ONE COUNTRY = ONE PROFILE", 11.5, True, ACCENT)]},
         {"runs": [("Each of the 193 countries becomes a single 11-feature trajectory profile - "
                    "the clustering unit for everything that follows.", 12, False, INK)],
          "space_before": 6}])
tag(s, r3, r4)
r5 = rect(s, 7.35, 5.52, 5.1, 1.15, WHITE, line_color=ACCENT)
r6 = tb(s, 7.55, 5.64, 4.7, 0.95,
        [{"runs": [("AUDITABLE BY DESIGN", 11.5, True, ACCENT)]},
         {"runs": [("Every exclusion is logged with its reason in cohort_audit.csv, so the "
                    "sample can be reproduced or challenged.", 11.5, False, INK)],
          "space_before": 4}])
tag(s, r5, r6)
notes(s, "This funnel defines the modelling sample and makes every exclusion traceable: from "
         "6,378 country-year rows (213 entities) we remove 16 micro states and 4 entities without "
         "a complete 1992-2021 history, leaving 193 modelled countries. Design rationale: a "
         "narrowing funnel encodes attrition visually; exclusion reasons sit next to the counts "
         "so the audience can audit the sample. Source: cohort_audit.csv from our pipeline. " + GENAI)

# ---- features
s = new_slide(STAGE_DESIGN, "Eleven features describe level, direction, timing and fuel shift",
              "2 / 4", "RAN CHANGMING", 12,
              topic="Turning 30-year emission histories into model features")
cards = [
    ("01 · LEVEL", "How intense is the economy now?",
     ["Log recent per-capita CO2", "(median 2017-2021)"],
     "Log scale stops extreme values from dominating distances.", LBLUE, BLUE),
    ("02 · DIRECTION", "Which way is the trajectory moving?",
     ["Long-term log slope 1992-2021", "Recent log slope 2012-2021",
      "Acceleration (recent - long-term)"],
     "Separates long-run direction from recent slowdown or renewed growth.", LORANGE, ORANGE),
    ("03 · TIMING", "When does change happen?",
     ["Growth volatility (MAD)", "Peak timing (relative)", "Post-peak change"],
     "Distinguishes smooth plateauing from unstable or still-rising paths.", LBLUE, BLUE),
    ("04 · FUEL SHIFT", "Is the fuel mix decarbonising?",
     ["Recent coal & gas shares", "Coal & gas share change", "(early vs recent 5-yr medians)"],
     "Oil share is kept for interpretation but omitted from the model to avoid perfect dependence.",
     LORANGE, ORANGE),
]
xs = [0.5, 3.62, 6.74, 9.86]
for (head, q, feats, why, fill, accent), x in zip(cards, xs):
    r1 = rect(s, x, 1.25, 2.9, 4.4, fill)
    r2 = rect(s, x, 1.25, 2.9, 0.09, accent)
    r3 = tb(s, x + 0.16, 1.46, 2.6, 0.3, [{"runs": [(head, 11.5, True, accent)]}])
    r4 = tb(s, x + 0.16, 1.80, 2.6, 0.8, [{"runs": [(q, 13, True, INK)]}])
    r5 = tb(s, x + 0.16, 2.68, 2.6, 1.8,
            [{"runs": [("· " + f, 12.5, False, INK)], "space_before": 4} for f in feats])
    r6 = tb(s, x + 0.16, 4.58, 2.6, 1.0, [{"runs": [(why, 11.5, False, MIDGREY)]}])
    tag(s, r1, r2, r3, r4, r5, r6)
b1 = rect(s, 0.5, 5.85, 11.95, 0.82, GREY)
b2 = tb(s, 0.72, 5.96, 11.5, 0.65,
        [{"runs": [("Built only from per-capita CO2, growth rates and fuel shares - ", 12.5, True, INK),
                   ("total CO2 and population never enter the model; they return later, "
                    "only to interpret the pathways.", 12.5, False, ACCENT)]}])
tag(s, b1, b2)
notes(s, "Feature engineering turns each 30-year history into 11 comparable features in four "
         "blocks: level (log recent per-capita CO2), direction (long-term slope, recent slope, "
         "acceleration), timing (volatility, peak timing, post-peak change) and fuel shift "
         "(coal/gas shares and their change). Deliberate exclusion: total emissions and population "
         "are NOT features - size must not decide the grouping; size returns only when we interpret "
         "the pathways in Visual Storytelling. Design rationale: four parallel cards mirror the "
         "four feature families so the audience learns the feature grammar once. " + GENAI)

# ---- pipeline
s = new_slide(STAGE_DESIGN, "A fixed pipeline turns 30-year histories into two pathways",
              "3 / 4", "RAN CHANGMING", 13, refs="[4]",
              topic="Scaling, PCA, K-means and how they are validated")
steps = [
    ("COUNTRY-YEAR PANEL", "6,378 rows · 1992-2021", LBLUE),
    ("193 COUNTRY PROFILES", "one row per country, 11 features", LBLUE),
    ("ROBUSTSCALER", "median / IQR scaling, so outliers cannot dominate distances", GREY),
    ("PCA \u2192 4 COMPONENTS", "87.41% of variance kept, correlated features merged", GREY),
    ("K-MEANS · K = 2-8", "50 restarts, fixed seed 42: the grouping re-runs identically",
     LORANGE),
]
x = 0.5
step_shapes = []
for i, (head, sub, fill) in enumerate(steps):
    r1 = rect(s, x, 1.5, 2.12, 1.5, fill, line_color=LINEGREY)
    r2 = tb(s, x + 0.08, 1.6, 1.96, 0.6,
            [{"runs": [(head, 11.5, True, INK)], "align": PP_ALIGN.CENTER}],
            align=PP_ALIGN.CENTER)
    r3 = tb(s, x + 0.08, 2.2, 1.96, 0.74,
            [{"runs": [(sub, 11, False, MIDGREY)], "align": PP_ALIGN.CENTER}],
            align=PP_ALIGN.CENTER)
    shapes = [r1, r2, r3]
    if i < 4:
        shapes.append(rect(s, x + 2.16, 2.13, 0.32, 0.24, ACCENT,
                           shape=MSO_SHAPE.RIGHT_ARROW))
    tag(s, *shapes)
    x += 2.5
label(s, 0.5, 3.15, 11.95, "VALIDATION RUNS ALONGSIDE - NOT INSIDE - THE CLUSTERING")
v1 = rect(s, 0.5, 3.5, 3.85, 1.2, WHITE, line_color=ORANGE)
v2 = tb(s, 0.68, 3.62, 3.55, 1.0,
        [{"runs": [("GMM SOFT MEMBERSHIP", 11.5, True, ORANGE)]},
         {"runs": [("probability of belonging to each cluster", 11, False, INK)],
          "space_before": 2}])
v3 = rect(s, 4.55, 3.5, 3.85, 1.2, WHITE, line_color=ORANGE)
v4 = tb(s, 4.73, 3.62, 3.55, 1.0,
        [{"runs": [("100 x 80% SUBSAMPLING", 11.5, True, ORANGE)]},
         {"runs": [("refit and compare labels (Adjusted Rand Index)", 11, False, INK)],
          "space_before": 2}])
v5 = rect(s, 8.6, 3.5, 3.85, 1.2, GREY)
v6 = tb(s, 8.78, 3.62, 3.55, 1.0,
        [{"runs": [("NEITHER CHANGES THE LABELS", 11.5, True, ACCENT)]},
         {"runs": [("both instruments only test the grouping that K-means already produced",
                    11, False, INK)], "space_before": 2}])
tag(s, v1, v2, v3, v4)
tag(s, v5, v6)
label(s, 0.5, 4.90, 6.0, "WHY K-MEANS")
w1 = rect(s, 0.5, 5.25, 11.95, 1.4, GREY)
w2 = tb(s, 0.72, 5.36, 11.5, 1.2,
        [{"runs": [("Unsupervised by design: ", 12.5, True, INK),
                   ("no labels exist for 'transition pathway', so we cluster comparable profiles "
                    "and validate the structure afterwards.", 12.5, False, INK)]},
         {"runs": [("Simple, explainable, reproducible: ", 12.5, True, INK),
                   ("a fixed seed and n_init = 50 make the exact grouping re-runnable; PCA removes "
                    "redundancy among correlated trajectory features before distances are computed.",
                    12.5, False, INK)], "space_before": 5}])
tag(s, w1, w2)
notes(s, "The pipeline is deliberately fixed end-to-end so every later figure traces back to one "
         "model run: panel -> 193 eleven-feature profiles -> RobustScaler -> PCA (4 components, "
         "87.41% cumulative variance) -> K-means compared over K = 2-8 (n_init = 50, seed = 42). "
         "GMM soft membership and 100 x 80% subsampling are validation instruments drawn outside "
         "the main flow. Design rationale: a left-to-right pipeline matches reading order; the "
         "amber validation boxes hang below to show they check the model rather than train it. " + GENAI)

# ---- K selection
s = new_slide(STAGE_DESIGN, "All three validity metrics rank K = 2 first",
              "4 / 4", "RAN CHANGMING", 14, refs="[4] [5]",
              topic="Choosing how many emission patterns to keep (K = 2-8)")
pic = picture_fit(s, FIG / "fig_k_selection.png", 0.32, 1.08, 9.5, 3.0)
cap = caption(s, 0.5, 4.14, 9.3,
              "Inertia (within-cluster sum of squares) always falls as K grows - shown for "
              "reference only, never used alone to pick K.")
tag(s, pic, cap)
w1 = rect(s, 0.5, 4.58, 9.3, 1.98, GREY)
w2 = tb(s, 0.72, 4.70, 8.9, 1.8,
        [{"runs": [("How the decision was made", 12, True, ACCENT)]},
         {"runs": [("· Three internal criteria were computed on the same PCA space for every "
                    "candidate K from 2 to 8, so the comparison is like-for-like.",
                    12, False, INK)], "space_before": 6},
         {"runs": [("· All three agree on K = 2; no criterion prefers a finer split, and "
                    "K = 3-8 lose separation without adding interpretable structure.",
                    12, False, INK)], "space_before": 4},
         {"runs": [("· K = 2 is therefore a screening frame for monitoring attention, not a "
                    "claim that only two kinds of country exist.", 12, False, INK)],
          "space_before": 4}])
tag(s, w1, w2)
stats = [("0.486", "Silhouette", "highest - compact, well separated"),
         ("0.865", "Davies-Bouldin", "lowest - least overlap"),
         ("98.3", "Calinski-Harabasz", "highest - strongest split")]
y = 1.3
for big, lab_, sub in stats:
    r1 = rect(s, 9.95, y, 2.85, 1.08, WHITE, line_color=BLUE)
    r2 = tb(s, 9.95, y + 0.06, 2.85, 0.44,
            [{"runs": [(big, 19, True, BLUE), ("   " + lab_, 12, True, INK)]}])
    r3 = tb(s, 9.95, y + 0.56, 2.85, 0.42, [{"runs": [(sub, 11.5, False, MIDGREY)]}])
    tag(s, r1, r2, r3)
    y += 1.24
d1 = rect(s, 9.95, y + 0.06, 2.85, 1.75, LORANGE)
d2 = tb(s, 10.11, y + 0.2, 2.55, 1.55,
        [{"runs": [("DECISION", 11.5, True, ORANGE)]},
         {"runs": [("K = 2 - the strongest macro segmentation, not the most detailed taxonomy.",
                    12.5, True, INK)], "space_before": 5}])
tag(s, d1, d2)
notes(s, "K was chosen by evidence, not convenience: silhouette (0.486, highest), Davies-Bouldin "
         "(0.865, lowest) and Calinski-Harabasz (98.3, highest) all rank K = 2 first across "
         "K = 2-8. Inertia is plotted for completeness only because it mechanically decreases "
         "with K. Design rationale: three small-multiple panels share one x-axis so the same "
         "dashed marker at K = 2 can be compared across criteria at a glance. "
         "Source: k_selection_metrics.csv. " + GENAI)

# =================================================================
# BODY 13-15 — MODEL EVALUATION AND PERFORMANCE
# =================================================================
s = new_slide(STAGE_EVAL, "Two pathways separate cleanly in reduced feature space",
              "1 / 3", "REN SIYU", 15,
              topic="Does the two-pattern structure actually separate?")
pic = picture_fit(s, FIG / "fig_pca_pathways.png", 0.28, 1.04, 8.8, 5.6)
tag(s, pic)
label(s, 9.2, 1.22, 3.9, "HOW TO READ")
h1 = tb(s, 9.2, 1.56, 3.3, 2.9,
        [{"runs": [("Each dot is one country's 11-feature profile, projected onto the first two "
                    "principal components.", 12, False, INK)]},
         {"runs": [("Colour shows the K-means label computed in the full 4-component space - "
                    "not in this 2D view.", 12, False, INK)], "space_before": 8},
         {"runs": [("Dark rings mark the 10 low-confidence boundary countries.", 12, False, INK)],
          "space_before": 8}])
tag(s, h1)
r1 = rect(s, 9.2, 4.55, 3.3, 1.95, GREY)
r2 = tb(s, 9.36, 4.68, 3.0, 1.75,
        [{"runs": [("PC1 + PC2 show 70.8% of variance", 13, True, ACCENT)]},
         {"runs": [("Clustering used all four retained components (87.41%). This 2D projection is "
                    "a display device - the visible separation understates the true separation.",
                    11.5, False, INK)], "space_before": 5}])
tag(s, r1, r2)
notes(s, "A PCA projection checks whether the K = 2 structure is real in the feature space rather "
         "than an artefact of the algorithm. The amber rapid-growth pathway (24 countries) forms "
         "a compact group clearly offset from the broader steel-blue plateauing/slowdown pathway "
         "(169). Dark rings flag the 10 boundary countries whose membership is uncertain. Caveat "
         "stated on-slide: clustering ran on 4 components (87.41% variance); PC1+PC2 alone show "
         "70.8%, so this view is illustrative. Design rationale: colour encodes only pathway "
         "membership; a third visual channel (ring) carries uncertainty instead of a third "
         "colour. " + GENAI)

# ---- internal metrics table
s = new_slide(STAGE_EVAL, "K = 2 wins on every internal validity metric",
              "2 / 3", "REN SIYU", 16, refs="[5]",
              topic="Internal validity metrics, since no ground-truth labels exist")
headers = ["K", "Inertia (lower)", "Silhouette (higher)", "Davies-Bouldin (lower)",
           "Calinski-Harabasz (higher)"]
data = [
    ("2", "1675.8", "0.486", "0.865", "98.3"),
    ("3", "1296.4", "0.361", "0.969", "91.0"),
    ("4", "1064.7", "0.271", "1.168", "87.2"),
    ("5", "915.9", "0.278", "1.086", "83.2"),
    ("6", "782.4", "0.283", "1.004", "83.9"),
    ("7", "700.6", "0.265", "1.071", "81.3"),
    ("8", "647.4", "0.263", "1.045", "77.2"),
]
tbl_shape = s.shapes.add_table(8, 5, Inches(0.5), Inches(1.3), Inches(7.05), Inches(4.5))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(0.7)
for j in range(1, 5):
    tbl.columns[j].width = Inches(1.5875)
for j, h in enumerate(headers):
    c = tbl.cell(0, j)
    c.text = h
    c.fill.solid()
    c.fill.fore_color.rgb = INK
    for p in c.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.color.rgb = WHITE
for i, row in enumerate(data):
    for j, v in enumerate(row):
        c = tbl.cell(i + 1, j)
        c.text = v
        sel = (i == 0)
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT if sel else (GREY if i % 2 else WHITE)
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(12.5)
                r.font.bold = sel
                r.font.color.rgb = WHITE if sel else INK
tb(s, 0.5, 5.9, 7.05, 0.5,
   [{"runs": [("Highlighted row: K = 2 wins Silhouette, Davies-Bouldin and Calinski-Harabasz "
               "simultaneously.", 11.5, True, ACCENT)]}])
meaning = [
    ("SILHOUETTE", "Rewards clusters that are internally compact and well separated from neighbours."),
    ("DAVIES-BOULDIN", "Penalises overlap between cluster pairs; lower means cleaner separation."),
    ("CALINSKI-HARABASZ", "Compares between-cluster with within-cluster dispersion; higher is better."),
]
y = 1.3
for head, txt in meaning:
    r1 = rect(s, 7.85, y, 4.6, 0.95, GREY)
    r2 = tb(s, 8.03, y + 0.08, 4.3, 0.82,
            [{"runs": [(head, 11.5, True, BLUE)]},
             {"runs": [(txt, 11, False, INK)]}])
    tag(s, r1, r2)
    y += 1.08
n1 = rect(s, 7.85, y + 0.06, 4.6, 1.6, WHITE, line_color=ACCENT)
n2 = tb(s, 8.03, y + 0.2, 4.3, 1.4,
        [{"runs": [("NO SUPERVISED ACCURACY", 11.5, True, ACCENT)]},
         {"runs": [("Accuracy, precision, recall, F1, ROC and confusion matrices require "
                    "ground-truth labels. Pathways have none - evaluation therefore uses internal "
                    "validity plus stability.", 11, False, INK)], "space_before": 4}])
tag(s, n1, n2)
notes(s, "The full K = 2-8 comparison table lets the audience verify the K choice themselves; "
         "each column header states the direction of 'better'. We explicitly state why "
         "supervised metrics are not applicable - clustering has no ground-truth labels, so "
         "internal validity (this slide) plus stability and membership uncertainty (next slide) "
         "carry the evaluation. Design rationale: the selected row uses the deck's maroon accent to "
         "pull the eye to the decision before the supporting rows are read. "
         "Source: k_selection_metrics.csv. " + GENAI)

# ---- stability + boundary
s = new_slide(STAGE_EVAL, "The split survives resampling; ten stay uncertain",
              "3 / 3", "REN SIYU", 17, refs="[6]",
              topic="Resampling stability and membership confidence")
pic = picture_fit(s, FIG / "fig_stability_gmm.png", 0.28, 1.04, 8.95, 3.9)
cap = caption(s, 0.5, 5.06, 8.8,
              "Left: 100 refits on 80% country subsamples - median Adjusted Rand Index 0.965 "
              "(minimum 0.778). Right: GMM soft membership on the same feature space.")
tag(s, pic, cap)
label(s, 9.4, 1.22, 3.5, "BOUNDARY COUNTRIES")
b1 = rect(s, 9.4, 1.56, 3.05, 3.1, LORANGE)
b2 = tb(s, 9.58, 1.68, 2.7, 2.9,
        [{"runs": [("GMM confidence below 0.70:", 11.5, True, ORANGE)]},
         {"runs": [("South Korea 0.51 · Guatemala 0.53 · Trinidad & Tobago 0.56 · Sri Lanka 0.61 "
                    "· Estonia 0.63 · Venezuela 0.65 · Oman 0.66 · Jordan 0.66 · Madagascar 0.67 "
                    "· Azerbaijan 0.68", 11, False, INK)], "space_before": 4},
         {"runs": [("Not a third pathway - membership to be confirmed with more data.",
                    11, True, ACCENT)], "space_before": 6}])
tag(s, b1, b2)
v1 = rect(s, 0.5, 5.62, 11.95, 0.95, GREY)
v2 = tb(s, 0.72, 5.72, 11.5, 0.8,
        [{"runs": [("Verdict: ", 12.5, True, ACCENT),
                   ("the two-pathway macro structure is stable (ARI close to 1 on most refits), "
                    "and the only genuine uncertainty is a small, known boundary set - exactly "
                    "the countries a monitoring programme should track more closely.",
                    12.5, False, INK)]}])
tag(s, v1, v2)
notes(s, "Stability was tested with 100 refits on 80% country subsamples (scaling, PCA and "
         "K-means re-estimated each time): median ARI 0.965 against the full-sample labels, so the "
         "macro split is not an initialisation artefact. A GMM fitted on the same PCA space adds "
         "soft membership probabilities: median confidence 98.7%, and exactly 10 countries fall "
         "below 0.70 - these become the 'boundary countries' carried into Visual Storytelling. "
         "Design rationale: histogram plus sorted dot plot show the full distributions rather than "
         "a single summary statistic. Sources: stability_results.csv, "
         "gmm_membership_probabilities.csv. " + GENAI)

# =================================================================
# BODY 16-19 — VISUAL STORYTELLING  (every page is figure-led)
# =================================================================
# ---- 1/4: aggregate decomposition + per-capita trajectories
s = new_slide(STAGE_STORY, "One rising total conceals two different transition pressures",
              "1 / 4", "REN SIYU", 18, refs="[1] [2]",
              topic="What kind of emission patterns did the model find?")
pic = picture_fit(s, FIG / "fig_pathway_profiles.png", 0.28, 0.97, 12.8, 3.12)
cap = caption(s, 0.52, 4.10, 12.3,
              "Shares use 2017-2021 emissions of the 193 modelled countries; the 10 boundary "
              "countries stay inside their pathway, not a third group.")
tag(s, pic, cap)


def pathway_strip(x, fill, accent, count, name, thesis, metrics, fuel, countries):
    r1 = rect(s, x, 4.42, 5.95, 2.32, fill)
    r2 = rect(s, x, 4.42, 5.95, 0.09, accent)
    r3 = tb(s, x + 0.18, 4.58, 4.0, 0.42, [{"runs": [(name, 14.5, True, INK)]}])
    r4 = tb(s, x + 4.2, 4.53, 1.55, 0.5,
            [{"runs": [(count, 23, True, accent), (" countries", 11, False, MIDGREY)],
              "align": PP_ALIGN.RIGHT}], align=PP_ALIGN.RIGHT)
    r5 = tb(s, x + 0.18, 5.02, 5.6, 0.42, [{"runs": [(thesis, 11.5, True, accent)]}])
    tag(s, r1, r2, r3, r4, r5)
    mx = x + 0.18
    for big, lab_ in metrics:
        m1 = rect(s, mx, 5.50, 1.4, 0.86, WHITE, line_color=accent)
        m2 = tb(s, mx + 0.04, 5.53, 1.32, 0.32,
                [{"runs": [(big, 13, True, accent)], "align": PP_ALIGN.CENTER}],
                align=PP_ALIGN.CENTER)
        m3 = tb(s, mx + 0.04, 5.85, 1.32, 0.48,
                [{"runs": [(lab_, 10.5, False, INK)], "align": PP_ALIGN.CENTER}],
                align=PP_ALIGN.CENTER)
        tag(s, m1, m2, m3)
        mx += 1.44
    f1 = tb(s, x + 0.18, 6.40, 5.6, 0.34,
            [{"runs": [("Median country shares  ", 10.5, True, MIDGREY), (fuel, 10.5, True, INK),
                       ("   ", 10.5, False, INK), (countries, 10.5, False, INK)]}])
    tag(s, f1)


pathway_strip(0.5, LORANGE, ORANGE, "24", "Rapid-growth pathway",
              "Lower base, fast recent growth, late peak - potential future additional "
              "emissions if recent patterns persist.",
              [("1.49 t", "CO2 per person"), ("+2.78%", "recent growth per year"),
               ("2018", "median peak year"), ("7.8%", "share of 2017-21 emissions")],
              "coal 38% · gas 2%",
              "e.g. Vietnam, Bangladesh")
pathway_strip(6.5, LBLUE, BLUE, "169", "Plateauing / slowdown pathway",
              "Growth has flattened or turned down - yet this pathway still carries almost "
              "all of the current emissions burden.",
              [("3.19 t", "CO2 per person"), ("-0.26%", "recent growth per year"),
               ("2009", "median peak year"), ("92.2%", "share of 2017-21 emissions")],
              "coal 1% · gas 11%",
              "e.g. China, United States")
notes(s, "The stacked area on the left decomposes the rising sample total: the plateauing "
         "pathway is the thick base that carries 92.2% of 2017-2021 emissions, while the "
         "rapid-growth band is thin but grew far faster (+198% versus +58% since 1992). The "
         "right panel shows the two transition shapes: the plateauing pathway peaks around 2009 "
         "and flattens, the rapid-growth pathway is still climbing with a median peak only in "
         "2018. Pathways are empirical patterns in the sample window - not country rankings, not "
         "predictions. Fuel shares are the median country-level CO2-source shares inside each "
         "pathway - the typical country, not the pathway aggregate: weighted by emissions the two "
         "pathways look similar (coal 48% vs 44% of 2017-2021 fossil CO2), because a few large "
         "coal economies sit inside the plateauing pathway. Design "
         "rationale: the figure carries the argument and the two strips below only quantify it, so "
         "this page is evidence-led rather than text-led. Sources: co2_panel_1992_2021.csv, "
         "cluster_profiles.csv. " + GENAI)

# ---- 2/4: burden x growth
s = new_slide(STAGE_STORY, "Current burden and future pressure require different lenses",
              "2 / 4", "REYIZA RESIBIEKE", 19,
              topic="Why ranking countries by one number misleads")
pic = picture_fit(s, FIG / "fig_burden_growth.png", 0.28, 1.02, 9.35, 5.55)
tag(s, pic)
label(s, 9.68, 1.22, 3.0, "WHY ONE LENS FAILS")
r1 = rect(s, 9.68, 1.58, 2.8, 2.2, LORANGE)
r2 = tb(s, 9.86, 1.70, 2.45, 2.0,
        [{"runs": [("Rank by emissions only ", 11.5, True, ORANGE)]},
         {"runs": [("and you watch China, the US and India while missing small-base, "
                    "fast-growing, late-peak countries in the upper left.", 11, False, INK)]}])
tag(s, r1, r2)
r3 = rect(s, 9.68, 3.95, 2.8, 2.2, LBLUE)
r4 = tb(s, 9.86, 4.07, 2.45, 2.0,
        [{"runs": [("Rank by growth only ", 11.5, True, BLUE)]},
         {"runs": [("and you overrate small economies while underrating the plateauing majors "
                    "that hold 92.2% of 2017-21 emissions.", 11, False, INK)]}])
tag(s, r3, r4)
notes(s, "Each country is plotted by current burden (x: average annual fossil CO2 2017-2021, log "
         "scale) against recent momentum (y: 2012-2021 log-slope of per-capita CO2, % per year); "
         "colour is the validated pathway and dark rings are boundary countries. Reading: upper "
         "right = burden AND growth (double attention); upper left = future pressure; lower right "
         "= today's mitigation lever. Both single-lens failures are stated on the slide. Bubble "
         "area is deliberately NOT reused for emissions because the x-axis already carries it. "
         + GENAI)

# ---- 3/4: pathway map
s = new_slide(STAGE_STORY, "Transition pressures are geographically uneven",
              "3 / 4", "SHEN RUITING", 20, refs="[1] [3]",
              topic="Where each emission pattern is located")
pic = picture_fit(s, FIG / "fig_pathway_map.png", 0.30, 0.99, 12.75, 5.35)
tag(s, pic)
b2 = tb(s, 0.62, 6.36, 11.9, 0.44,
        [{"runs": [("The map locates monitoring tasks, not judgement: ", 13, True, ACCENT),
                   ("rapid-growth pressure clusters in South and Southeast Asia; dark "
                    "outlines flag the 10 boundary countries.", 13, False, INK)]}])
tag(s, b2)
notes(s, "The frozen K = 2 labels are projected onto the world map using only the two pathway "
         "colours; boundary countries get a dark outline instead of a third colour, and grey "
         "countries were not modelled (micro states or incomplete histories). Geography was never "
         "a model feature - this map shows where the data-driven pathways happen to fall. "
         "Twenty-eight small island or city-state members (27 plateauing, 1 rapid-growth) are "
         "invisible at this 1:110m base-map scale. Design rationale: the same basemap family and "
         "legend style as the exploration maps keeps the spatial language consistent. "
         "Source: country_features_with_clusters.csv + Natural Earth 110m. " + GENAI)

# ---- 4/4: centre of gravity + differentiated monitoring
s = new_slide(STAGE_STORY, "Where emissions head decides where monitoring goes",
              "4 / 4", "SHEN RUITING", 21, refs="[3]",
              topic="How limited monitoring attention should be allocated")
# the figure carries its own year legend, so no separate caption is needed here
pic = picture_fit(s, FIG / "fig_pathway_centroids.png", 0.26, 0.97, 8.95, 4.35)
tag(s, pic)
moved = [
    ("15E to 44E", "Centre of gravity of all 193 countries shifts east.", BLUE, LBLUE),
    ("+198%", "Rapid-growth emissions since 1992; centre stays in South/Southeast Asia.",
     ORANGE, LORANGE),
    ("+58%", "Plateauing emissions - 92.2% of 2017-21 emissions.", BLUE, LBLUE),
]
x = 0.5
for big, txt, accent, fill in moved:
    r1 = rect(s, x, 5.28, 2.63, 1.02, fill)
    r2 = tb(s, x + 0.16, 5.32, 2.3, 0.38, [{"runs": [(big, 15, True, accent)]}])
    r3 = tb(s, x + 0.16, 5.68, 2.3, 0.56, [{"runs": [(txt, 11, False, INK)]}])
    tag(s, r1, r2, r3)
    x += 2.72
# the storyline document requires this caveat on the monitoring page: the same
# symptom can have different causes, so K = 2 screens rather than diagnoses
foot = tb(s, 0.52, 6.36, 8.3, 0.34,
          [{"runs": [("Same symptom, different causes: ", 10.5, True, ACCENT),
                     ("flat emissions can mean expansion, slow post-peak decline or a real "
                      "fuel shift.", 10.5, False, MIDGREY)]}])
tag(s, foot)
label(s, 9.32, 1.10, 4.0, "MONITORING PRIORITIES")
mon = [
    ("Rapid-growth pathway (24)",
     "Watch: growth rate · fuel shares · peak timing",
     "Progress = growth slows, high-carbon shares fall",
     ORANGE, LORANGE),
    ("Plateauing pathway (169)",
     "Watch: total emissions · pace of decline",
     "Progress = absolute emissions fall structurally",
     BLUE, LBLUE),
    ("Boundary countries (10)",
     "Watch: data completeness · membership stability",
     "Progress = pathway membership stabilises with better data",
     MIDGREY, GREY),
]
y = 1.45
for name, watch, prog, accent, fill in mon:
    r1 = rect(s, 9.32, y, 3.35, 1.68, fill)
    r2 = rect(s, 9.32, y, 0.07, 1.68, accent)
    r3 = tb(s, 9.53, y + 0.06, 3.0, 0.4, [{"runs": [(name, 12.5, True, accent)]}])
    r4 = tb(s, 9.53, y + 0.46, 3.0, 1.15,
            [{"runs": [(watch, 11.5, False, INK)]},
             {"runs": [(prog, 11.5, False, INK)], "space_before": 4}])
    tag(s, r1, r2, r3, r4)
    y += 1.80
notes(s, "This closing storytelling page joins the spatial and the managerial argument. The map "
         "reuses the spatial team's centroid method (CO2-weighted country centroids, "
         "code/瑞庭_绘图/scripts/_render_p3_trajectory_revised.py) but replaces continents with the "
         "two validated pathways: the overall centre of gravity moves from 15E to 44E, the "
         "plateauing centre follows it, and the rapid-growth centre sits in South/Southeast Asia "
         "the whole time while its total grows +198%. The table then converts pathways into "
         "monitoring instructions: growth, fuel shares, peak timing and lock-in for the "
         "rapid-growth pathway; absolute decline for the plateauing pathway; data quality and "
         "membership stability for boundary countries. Positions are plotted at true coordinates "
         "with no magnification. " + GENAI)

# =================================================================
# BODY 20 — CONCLUSIONS
# =================================================================
s = new_slide("CONCLUSIONS", "From one emission total to pathway-specific monitoring",
              None, "GROUP 30 - ALL MEMBERS", 22,
              topic="What this investigation changes for climate monitoring")
cards = [
    ("CURRENT BURDEN", "The 169-country plateauing pathway holds 92.2% of 2017-2021 sample emissions",
     "Push the plateau into absolute decline - this is the short-term mitigation lever.", BLUE, LBLUE),
    ("FUTURE PRESSURE", "The 24-country rapid-growth pathway: fast growth, late peaks, +198% since 1992",
     "Track growth and lock-in signals now, before today's growth becomes tomorrow's burden.",
     ORANGE, LORANGE),
    ("BOUNDARY UNCERTAINTY", "10 low-confidence countries identified by soft membership",
     "Observe them more frequently instead of forcing a fixed classification.", MIDGREY, GREY),
]
xs = [0.5, 4.53, 8.56]
for (head, stat, act, accent, fill), x in zip(cards, xs):
    r1 = rect(s, x, 1.3, 3.9, 2.35, fill)
    r2 = rect(s, x, 1.3, 3.9, 0.1, accent)
    r3 = tb(s, x + 0.18, 1.48, 3.55, 0.35, [{"runs": [(head, 11.5, True, accent)]}])
    r4 = tb(s, x + 0.18, 1.86, 3.55, 0.88, [{"runs": [(stat, 12, True, INK)]}])
    r5 = tb(s, x + 0.18, 2.76, 3.55, 0.85, [{"runs": [(act, 11.5, False, INK)]}])
    tag(s, r1, r2, r3, r4, r5)
n1 = rect(s, 0.5, 3.92, 11.96, 1.06, WHITE, line_color=BLUE)
n2 = tb(s, 0.72, 4.02, 11.5, 0.92,
        [{"runs": [("What is novel here: ", 13, True, BLUE),
                   ("an 11-feature trajectory profile (level · direction · timing · fuel shift) "
                    "that makes 30 years of history comparable per country, and a "
                    "stability-checked K = 2 screen used as a monitoring tool - not as another "
                    "league table.", 13, False, INK)]}])
tag(s, n1, n2)
l1 = rect(s, 0.5, 5.12, 11.96, 1.06, WHITE, line_color=ACCENT)
l2 = tb(s, 0.72, 5.22, 11.5, 0.92,
        [{"runs": [("Honest limits: ", 13, True, ACCENT),
                   ("pathways are empirical patterns in 1992-2021, not causes, forecasts or "
                    "policy verdicts; with no policy or economic variables in the panel, the "
                    "model prioritises deeper investigation - it does not replace it.",
                    13, False, INK)]}])
tag(s, l1, l2)
# the storyline document's closing sentence, kept as the last thing on the page
rect(s, 5.5, 6.34, 2.33, 0.035, ACCENT_L, glass=False)
c1 = tb(s, 0.5, 6.44, 12.33, 0.4,
        [{"runs": [("The purpose is not to rank countries, but to identify where different "
                    "transition pressures require different monitoring priorities.",
                    12.5, True, ACCENT)], "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER)
tag(s, c1)
notes(s, "Conclusions avoid repeating earlier slides and instead summarise the contribution: one "
         "rising sample total was decomposed into two validated transition pressures plus an "
         "explicitly uncertain boundary set, each with its own monitoring implication. Novelty: "
         "the trajectory-profile feature design and the stability-checked use of K-means as a "
         "screening device for management attention. Limits are stated plainly: exploratory, "
         "non-causal, non-predictive. " + GENAI)

# =================================================================
# ACKNOWLEDGEMENTS  (not counted)
# =================================================================
s = prs.slides.add_slide(LAYOUT)
background(s)
clear_placeholders(s)
banner(s, None, "Acknowledgements")
tb(s, 0.85, 1.24, 7.4, 1.0,
   [{"runs": [("Thank you", 30, True, INK)]},
    {"runs": [("for your attention and questions.", 15, False, MIDGREY)], "space_before": 6}])
rect(s, 0.85, 2.42, 1.5, 0.045, ACCENT)
ack = [
    ("Course", "CA6002 AI UX & Data Visualisation Design Principles, NTU."),
    ("Open-source tools", "Python, pandas, scikit-learn, matplotlib, GeoPandas, python-pptx."),
    ("Generative AI", "ChatGPT (OpenAI) assisted with Python plotting code; all data "
                      "processing, results and interpretation are the group's own work."),
]
y = 2.62
for head, txt in ack:
    rect(s, 0.85, y, 0.07, 0.66, ACCENT)
    tb(s, 1.08, y, 6.6, 0.66,
       [{"runs": [(head, 12, True, ACCENT)]},
        {"runs": [(txt, 11, False, INK)]}])
    y += 0.72
picture_fit(s, FIG / "fig_cover_art.png", 8.35, 1.28, 4.45, 2.25)
tb(s, 8.35, 3.62, 4.45, 1.1,
   [{"runs": [("Group 30", 15, True, ACCENT)]},
    {"runs": [("RAN CHANGMING  ·  REN SIYU", 12, True, INK)], "space_before": 4},
    {"runs": [("REYIZA RESIBIEKE  ·  SHEN RUITING", 12, True, INK)]}])
label(s, 0.85, 4.78, 4.0, "REFERENCES")
refs_left = [
    ("[1]", "Ritchie, H., Rosado, P. & Roser, M. (2023). CO\u2082 and Greenhouse Gas Emissions. "
            "Our World in Data."),
    ("[2]", "Friedlingstein, P. et al. (2022). Global Carbon Budget 2022. Earth Syst. Sci. "
            "Data, 14(11), 4811-4900."),
    ("[3]", "Natural Earth. Admin 0 - Countries, 1:110m Cultural Vectors (public domain)."),
]
refs_right = [
    ("[4]", "Pedregosa, F. et al. (2011). Scikit-learn: Machine Learning in Python. JMLR, 12, "
            "2825-2830."),
    ("[5]", "Rousseeuw, P. J. (1987). Silhouettes: a graphical aid to cluster analysis. "
            "J. Comput. Appl. Math., 20, 53-65."),
    ("[6]", "Hubert, L. & Arabie, P. (1985). Comparing partitions. J. Classification, 2, "
            "193-218."),
]
for col_x, col in ((0.85, refs_left), (6.75, refs_right)):
    yy = 5.12
    for num, txt in col:
        tb(s, col_x, yy, 0.42, 0.5, [{"runs": [(num, 9, True, ACCENT)]}])
        tb(s, col_x + 0.42, yy, 5.15, 0.5, [{"runs": [(txt, 9, False, MIDGREY)]}])
        yy += 0.58
notes(s, "Acknowledgements slide. It credits the course, the "
         "open-source toolchain and the generative-AI assistance, and lists the six references "
         "cited on the slides: the two data sources, the basemap, the clustering library and the "
         "two cluster-validity papers behind the silhouette and ARI diagnostics. " + GENAI)

# =================================================================
# Re-style the seven imported Exploration of Dataset slides
# =================================================================
EXPLORE_TITLES = [
    ("Where the emission data comes from", "Two sources create one traceable emissions panel"),
    ("Building one comparable country-year panel",
     "A controlled left join keeps the OWID backbone intact"),
    ("How much did the covered emissions rise, 1992-2021?",
     "Total CO2 grew nearly five times faster than per person"),
    ("Is the rise one shared trajectory or many?",
     "Long-term growth differs from recent momentum"),
    ("Who carries the current emission burden?",
     "A small set of countries drives most 2021 emissions"),
    ("Emission scale versus emission intensity",
     "Total burden and per-capita intensity tell different stories"),
    ("What the emissions are actually made of",
     "Fuel dependence is geographically uneven"),
]
# a closing line per imported slide: why the page matters for the argument
EXPLORE_NOTES = [
    ("Why it matters:  ",
     "every later feature is computed from these same 6,378 country-year rows, so the "
     "panel has to be auditable first."),
    ("Why it matters:  ",
     "conditioning is where a study can silently bias itself, so no observed value is "
     "overwritten and every filled field stays flagged."),
    ("Why it matters:  ",
     "the aggregate tells us the total grew, not which countries produced the growth "
     "or how their paths differ."),
    ("Reading the puzzle:  ",
     "long-term growth and recent momentum disagree, so no single growth number can "
     "describe a 30-year trajectory."),
    ("Why it matters:  ",
     "monitoring built around today's biggest emitters alone will overlook the "
     "fast-growing rest."),
    ("Why it matters:  ",
     "scale and intensity rank countries differently, which is the first evidence that "
     "one ranking cannot serve as one diagnosis."),
    ("Why it matters:  ",
     "fuel composition is what makes a transition easy or locked in, so coal and gas "
     "shares enter the model as the fuel-shift block."),
]
# the two dataset pages carry the data-source references listed on the closing slide
EXPLORE_REFS = {0: "Data: [1] [2]", 1: "Data: [1] [2]"}
# who presents each imported exploration page, in imported order
EXPLORE_RESP = ["SHEN RUITING", "SHEN RUITING", "REYIZA RESIBIEKE", "REYIZA RESIBIEKE",
                "SHEN RUITING", "SHEN RUITING", "SHEN RUITING"]
# the imported slides were drawn with a much more saturated palette; map the
# strong hues onto the muted deck palette so the whole deck reads as one design
RECOLOUR = {
    "4472C4": "4F6E96", "2E7FA0": "4F6E96", "ED7D31": "BE7F4E",
    "D94816": "BE7F4E", "A87300": "8A6E42", "C69A1A": "A98A52",
    "70AD47": "6E8C74", "800000": "2F5248", "7C2530": "2F5248",
    "66583D": "6B6250",
    "EAF2F8": "ECEFF4", "D9EAF3": "E4E9F0", "F7FAFC": "FAFBFA",
    "FFF0E6": "F7F1EA", "FFF7D6": "F6F2E4", "EDF6E8": "EEF4EC",
    "F8E8E9": "EEF4EC", "EAF4F1": "EDF3EF", "F5F7F9": "FAFBFA",
    "D2D9E0": "D5DED6", "F3F4F6": "FAFBFA",
}


def strip_chrome(slide):
    """Remove banner/wash chrome from a previous build so the pass is re-runnable.

    The exploration source deck is re-derived from the finished deck, which means
    the imported slides may already carry a banner bar, its hairline, its title
    textbox and the background wash. Only full-width shapes in the banner band -
    or a full-width wash - qualify; slide content is much narrower.
    """
    for sh in list(slide.shapes):
        if sh.width is None or sh.top is None or sh.height is None:
            continue
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:   # keep the template picture
            continue
        full_width = sh.width > Inches(9.0)
        in_banner = sh.top < Inches(0.95) and sh.height < Inches(1.0)
        is_wash = sh.top <= Inches(0.02) and sh.height > Inches(5.0) and full_width
        if full_width and (in_banner or is_wash):
            sh._element.getparent().remove(sh._element)


def tag_bands(slide, bands=3):
    """Reveal an imported slide in a few vertical bands.

    The imported slides were not authored with animation in mind, so instead of
    guessing their logical blocks the content is revealed top-to-bottom in three
    click steps: enough pacing for narration without a click per shape.
    """
    content = [sh for sh in slide.shapes
               if sh.top is not None and sh.height is not None
               and sh.top > Inches(0.95)
               and not (sh.shape_type == MSO_SHAPE_TYPE.PICTURE
                        and sh.width > Inches(12.5))]
    if not content:
        return
    content.sort(key=lambda sh: (sh.top, sh.left or 0))
    lo = min(sh.top for sh in content)
    hi = max(sh.top + sh.height for sh in content)
    step = max((hi - lo) // bands, 1)
    for b in range(bands):
        grp = [sh for sh in content if lo + b * step <= sh.top < lo + (b + 1) * step]
        if grp:
            tag(slide, *grp)


def wash_over_template(slide):
    """Imported slides carry the opaque white template picture, so the pale-green
    wash is laid on top of it and stops above the footer rule."""
    wash = gradient_rect(slide, 0, 0, SW_IN, 6.86, WASH_A, WASH_B, angle=45)
    spTree = slide.shapes._spTree
    spTree.remove(wash._element)
    spTree.insert(3, wash._element)


def recolour(shapes):
    for sh in shapes:
        if sh.shape_type == 6:
            recolour(sh.shapes)
            continue
        try:
            if sh.fill.type == 1:
                hexv = str(sh.fill.fore_color.rgb).upper()
                if hexv in RECOLOUR:
                    hexv = RECOLOUR[hexv]
                    sh.fill.fore_color.rgb = RGBColor.from_string(hexv)
                # near-white panels become frosted glass, like the new slides,
                # so they sit on the pale-green wash instead of punching through it
                r, g, b = (int(hexv[i:i + 2], 16) for i in (0, 2, 4))
                if min(r, g, b) > 235:
                    set_alpha(sh)
        except Exception:
            pass
        try:
            if sh.line.fill.type == 1:
                hexv = str(sh.line.color.rgb).upper()
                if hexv in RECOLOUR:
                    sh.line.color.rgb = RGBColor.from_string(RECOLOUR[hexv])
        except Exception:
            pass
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    try:
                        hexv = str(r.font.color.rgb).upper()
                    except Exception:
                        continue
                    if hexv in RECOLOUR:
                        r.font.color.rgb = RGBColor.from_string(RECOLOUR[hexv])


for i in range(7):
    slide = prs.slides[i]
    strip_chrome(slide)
    recolour(slide.shapes)
    wash_over_template(slide)
    # drop the stale kicker textboxes, the old plain-text title, and the imported
    # credit line, which used a different case, colour and position from ours
    for sh in list(slide.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip().upper().startswith(
                ("EXPLORATION OF DATASET", "MODEL EVALUATION", "VISUAL STORYTELLING",
                 "RESPONSIBLE:")):
            sh._element.getparent().remove(sh._element)
    if slide.shapes.title is not None:
        slide.shapes.title._element.getparent().remove(slide.shapes.title._element)
    topic_i, title_i = EXPLORE_TITLES[i]
    banner(slide, STAGE_EXPLORE, title_i, f"{i + 1} / 7", topic_i)
    # one line of narration in the empty band above the footer: it states what the
    # page contributes to the argument, which the original slides left implicit
    note = tb(slide, 0.85, 6.36, 11.6, 0.42,
              [{"runs": [(EXPLORE_NOTES[i][0], 11.5, True, ACCENT),
                         (EXPLORE_NOTES[i][1], 11.5, False, INK)]}])
    # credit line and page number in the same style as every other slide
    responsible(slide, EXPLORE_RESP[i])
    page_number(slide, i + 4)
    # the two dataset-introduction pages carry the source references
    if i in EXPLORE_REFS:
        ref_mark(slide, EXPLORE_REFS[i])
    tag_bands(slide)
    tag(slide, note)

# =================================================================
# Order: cover, contents, intro, exploration x7, design x4, eval x3, story x4,
#        conclusions, acknowledgements
# appended: [0-6 imported] [7 cover] [8 toc] [9 intro] [10-13 design]
#           [14-16 eval] [17-20 story] [21 conclusions] [22 acknowledgements]
# =================================================================
order = [7, 8, 9, 0, 1, 2, 3, 4, 5, 6, 10, 11, 12, 13, 14, 15, 16, 17, 18, 19, 20, 21, 22]
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
for el in ids:
    sldIdLst.remove(el)
for idx in order:
    sldIdLst.append(ids[idx])


# =================================================================
# Global typography pass: Times New Roman everywhere, no text below 11 pt
# =================================================================
def set_theme_fonts():
    """Point the theme's major/minor Latin typefaces at Times New Roman."""
    from lxml import etree
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    RT_THEME = ("http://schemas.openxmlformats.org/officeDocument/2006/"
                "relationships/theme")
    for master in prs.slide_masters:
        try:
            theme = master.part.part_related_by(RT_THEME)
        except KeyError:
            continue
        root = etree.fromstring(theme.blob)
        changed = False
        for name in ("majorFont", "minorFont"):
            el = root.find(f".//{A}{name}")
            if el is None:
                continue
            latin = el.find(f"{A}latin")
            if latin is not None:
                latin.set("typeface", FONT)
                changed = True
        if changed:
            theme._blob = etree.tostring(root, xml_declaration=True,
                                         encoding="UTF-8", standalone=True)


def restyle_text(shapes):
    for sh in shapes:
        if sh.shape_type == 6:  # group
            restyle_text(sh.shapes)
            continue
        if sh.has_text_frame:
            # narrow boxes: a short token ("Q1", "+") can still be set at body
            # size, but a phrase must be capped or it wraps into fragments
            words = sh.text_frame.text.split()
            short_token = words and max(len(w) for w in words) <= 4 and \
                len(sh.text_frame.text.strip()) <= 8
            # a single word in a narrow chip must never be split across lines
            if (sh.width is not None and sh.width < Inches(1.6) and len(words) == 1
                    and len(words[0]) > 2):
                sh.text_frame.word_wrap = False
            if sh.width is not None and sh.width < Inches(0.95) and not short_token:
                # a compact label ("SUPPLEMENT", "low", "PER PERSON") reads far
                # better spilling slightly than broken across lines
                if len(sh.text_frame.text.strip()) <= 12:
                    sh.text_frame.word_wrap = False
                for p in sh.text_frame.paragraphs:
                    if p.font.size is not None and p.font.size.pt > 9.0:
                        p.font.size = Pt(9.0)
                    if p.font.name != FONT_TITLE:
                        p.font.name = FONT
                    for r in p.runs:
                        if r.font.size is not None and r.font.size.pt > 9.0:
                            r.font.size = Pt(9.0)
                        if r.font.name != FONT_TITLE:
                            r.font.name = FONT
                continue
            tight = (sh.width is not None and sh.width < Inches(1.75)) or \
                    (sh.height is not None and sh.height < Inches(0.3))
            floor = 9.0 if tight else MIN_PT
            for p in sh.text_frame.paragraphs:
                if p.font.size is not None and p.font.size.pt < floor:
                    p.font.size = Pt(floor)
                if p.font.name != FONT_TITLE:
                    p.font.name = FONT
                for r in p.runs:
                    if r.font.size is not None and r.font.size.pt < floor:
                        r.font.size = Pt(floor)
                    if r.font.name != FONT_TITLE:
                        r.font.name = FONT
        if getattr(sh, "has_table", False) and sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        if p.font.size is not None and p.font.size.pt < MIN_PT:
                            p.font.size = Pt(MIN_PT)
                        p.font.name = FONT
                        for r in p.runs:
                            if r.font.size is not None and r.font.size.pt < MIN_PT:
                                r.font.size = Pt(MIN_PT)
                            r.font.name = FONT


def force_font_xml():
    """Catch the font references python-pptx cannot reach (a:ea, a:cs, list styles).

    Georgia is left alone: it is the deliberate display typeface for banners and
    cover headings.
    """
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for slide in prs.slides:
        for name in ("latin", "ea", "cs"):
            for el in slide._element.iter(f"{A}{name}"):
                if el.get("typeface") and el.get("typeface") != FONT_TITLE:
                    el.set("typeface", FONT)


set_theme_fonts()
for slide in prs.slides:
    restyle_text(slide.shapes)
force_font_xml()

inject_motion()
prs.save(str(OUT))
print("saved:", OUT, "slides:", len(list(prs.slides)))
