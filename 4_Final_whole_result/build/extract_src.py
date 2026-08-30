# -*- coding: utf-8 -*-
"""Rebuild the exploration-only source deck from the current final deck.

The original exploration deck was lost; its 7 slides survive inside the final deck,
so we copy the final deck and drop every other slide. The result is a build input,
written to build/source/ - never to the delivery folder. Also dumps embedded media
so the NTU logo can be reused.

Only needed if build/source/exploration_slides_source.pptx is ever lost again.
"""
import shutil
import zipfile
from pathlib import Path
from pptx import Presentation

ROOT = Path(r"C:\Users\user\Desktop\Overall_Data Visual Assignment")
FINAL = ROOT / "4_Final_whole_result" / "CA6002_Group30_Final_Presentation.pptx"
SRC = (ROOT / "4_Final_whole_result" / "build" / "source"
       / "exploration_slides_source.pptx")
SRC.parent.mkdir(parents=True, exist_ok=True)
ASSETS = ROOT / "4_Final_whole_result" / "build" / "assets"
ASSETS.mkdir(parents=True, exist_ok=True)

with zipfile.ZipFile(FINAL) as z:
    for n in z.namelist():
        if n.startswith("ppt/media/"):
            out = ASSETS / Path(n).name
            out.write_bytes(z.read(n))
            print("media:", out.name, out.stat().st_size)

shutil.copy(FINAL, SRC)
prs = Presentation(str(SRC))
keep = set(range(3, 10))  # 0-based: cover,toc,intro = 0-2; exploration = 3-9
sldIdLst = prs.slides._sldIdLst
for i, el in enumerate(list(sldIdLst)):
    if i not in keep:
        rId = el.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        sldIdLst.remove(el)
prs.save(str(SRC))
print("exploration source slides:", len(Presentation(str(SRC)).slides._sldIdLst))
for i, s in enumerate(Presentation(str(SRC)).slides, 1):
    t = s.shapes.title.text if s.shapes.title is not None else "(no title)"
    print(i, "|", t)
