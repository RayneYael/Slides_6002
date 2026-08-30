"""Remove the template's auto slide-number placeholder from the imported slides.

Those slides carry both our own page number and the layout's SLIDE_NUMBER
placeholder, which PowerPoint fills in by itself - hence two numbers on one
corner. The placeholder is deleted together with any animation entry that
targeted it, so the slide timing stays internally consistent.
"""
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

DECK = Path(__file__).resolve().parents[1] / "CA6002_Group30_Final_Presentation.pptx"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


def drop_timing_for(slide, spid):
    """Remove effect groups whose only target is the shape being deleted."""
    timing = slide._element.find(f"{P}timing")
    if timing is None:
        return 0
    gone = 0
    for par in list(timing.iter(f"{P}par")):
        tgts = {t.get("spid") for t in par.iter(f"{P}spTgt")}
        if tgts == {str(spid)} and par.getparent() is not None:
            par.getparent().remove(par)
            gone += 1
    return gone


prs = Presentation(str(DECK))
removed = effects = 0

for n, slide in enumerate(prs.slides, 1):
    for sh in list(slide.shapes):
        if not sh.is_placeholder:
            continue
        if sh.placeholder_format.type != PP_PLACEHOLDER.SLIDE_NUMBER:
            continue
        effects += drop_timing_for(slide, sh.shape_id)
        sh._element.getparent().remove(sh._element)
        removed += 1
        print(f"  slide {n}: removed auto slide-number placeholder")

prs.save(str(DECK))
print(f"placeholders removed: {removed}   stale animation groups removed: {effects}")
