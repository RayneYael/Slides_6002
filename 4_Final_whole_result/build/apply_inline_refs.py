"""Move the reference markers off the footer and into the body text.

Runs against the finished deck in place, so hand edits made in PowerPoint
survive: it only deletes the footer marker boxes and appends a small bracketed
number to the sentence or card that actually relies on each source. The
reference list on the closing slide is left untouched.
"""
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

DECK = Path(__file__).resolve().parents[1] / "CA6002_Group30_Final_Presentation.pptx"
MIDGREY = RGBColor(0x6E, 0x76, 0x72)
FOOTER_MARK = re.compile(r"^(Data:\s*)?(\[\d\]\s*)+$")

# slide number -> [(snippet that identifies the anchor text, marker)]
ANCHORS = {
    4:  [("OWID CO2 Data", "[1]"),
         ("Global Carbon Budget 2022", "[2]")],
    5:  [("valid ISO rows", "[1]"),
         ("fuel-only fill", "[2]")],
    13: [("50 restarts, fixed seed 42", "[4]")],
    14: [("Three internal criteria were computed", "[4] [5]")],
    16: [("Rewards clusters that are internally compact", "[5]")],
    17: [("median Adjusted Rand Index", "[6]")],
    18: [("Shares use 2017-2021 emissions", "[1] [2]")],
    20: [("The map locates monitoring tasks", "[1] [3]")],
    21: [("Centre of gravity of all 193 countries", "[3]")],
}


def append_marker(shape, snippet, marker):
    """Add the marker to the end of the paragraph that contains `snippet`."""
    for para in shape.text_frame.paragraphs:
        if snippet.lower() not in "".join(r.text for r in para.runs).lower():
            continue
        if marker in "".join(r.text for r in para.runs):
            return True
        last = para.runs[-1]
        run = para.add_run()
        run.text = f" {marker}"
        run.font.name = last.font.name or "Times New Roman"
        run.font.size = last.font.size or Pt(11)
        run.font.bold = False
        run.font.color.rgb = MIDGREY
        return True
    return False


prs = Presentation(str(DECK))
removed = placed = 0

for n, slide in enumerate(prs.slides, 1):
    for sh in list(slide.shapes):
        if (sh.has_text_frame and sh.top is not None and sh.top > Inches(6.8)
                and FOOTER_MARK.match(sh.text_frame.text.strip())):
            sh._element.getparent().remove(sh._element)
            removed += 1
    for snippet, marker in ANCHORS.get(n, []):
        hit = False
        for sh in slide.shapes:
            if sh.has_text_frame and append_marker(sh, snippet, marker):
                hit = True
                placed += 1
                break
        if not hit:
            print(f"  ! slide {n}: no anchor for {snippet!r}")

prs.save(str(DECK))
print(f"footer markers removed: {removed}   inline markers placed: {placed}")
