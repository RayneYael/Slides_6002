"""Check where reference markers now sit: inline in the body, none in the footer."""
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

DECK = Path(__file__).resolve().parents[1] / "CA6002_Group30_Final_Presentation.pptx"
prs = Presentation(str(DECK))
MARK = re.compile(r"\[\d\]")

for i, s in enumerate(prs.slides, 1):
    for sh in s.shapes:
        if not sh.has_text_frame or sh.top is None:
            continue
        for para in sh.text_frame.paragraphs:
            t = "".join(r.text for r in para.runs).strip()
            if not MARK.search(t):
                continue
            where = "FOOTER" if sh.top > Inches(6.8) else "body"
            print(f"slide {i:>2} [{where}] y={sh.top / 914400:.2f}  {t[:120]}")
