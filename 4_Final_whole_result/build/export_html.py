"""
Export the final PPTX to one self-contained, fully editable HTML file.

Every shape becomes an absolutely positioned <div> at the same coordinates as in
PowerPoint, so the HTML is a faithful copy of the deck rather than a re-design.
Text is contenteditable, shapes can be dragged, resized and deleted, pictures can
be swapped, and the edited document can be saved back to a new HTML file.

    python build/export_html.py

Output: CA6002_Group30_Final_Presentation.html  (images inlined as base64)
"""
from __future__ import annotations

import base64
import html
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).resolve().parent
DECK = HERE.parent / "CA6002_Group30_Final_Presentation.pptx"
OUT = HERE.parent / "CA6002_Group30_Final_Presentation.html"

PX = 96.0  # CSS pixels per inch; 13.333 x 7.5 in -> 1280 x 720 px
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}
ALIGN = {1: "left", 2: "center", 3: "right", 4: "justify"}
ANCHOR = {1: "flex-start", 2: "center", 3: "flex-end"}


def px(emu) -> float:
    return round(Emu(int(emu)).inches * PX, 2)


def hex_of(color_elm) -> str | None:
    """Read a colour from an a:solidFill / a:gs child element."""
    srgb = color_elm.find("a:srgbClr", NS)
    if srgb is not None:
        return "#" + srgb.get("val")
    scheme = color_elm.find("a:schemeClr", NS)
    if scheme is not None:
        return {"bg1": "#FFFFFF", "lt1": "#FFFFFF", "tx1": "#000000",
                "dk1": "#000000"}.get(scheme.get("val"), None)
    return None


def alpha_of(color_elm) -> float:
    node = color_elm.find(".//a:alpha", NS)
    if node is None:
        return 1.0
    return int(node.get("val")) / 100000.0


def rgba(hexcol: str, alpha: float) -> str:
    hexcol = hexcol.lstrip("#")
    r, g, b = (int(hexcol[i:i + 2], 16) for i in (0, 2, 4))
    return f"rgba({r},{g},{b},{round(alpha, 3)})"


def fill_css(shape) -> str:
    """Solid, gradient or no fill, read straight from the shape XML."""
    spPr = shape._element.find(".//a:spPr", NS)
    if spPr is None:
        return ""
    if spPr.find("a:noFill", NS) is not None:
        return ""
    solid = spPr.find("a:solidFill", NS)
    if solid is not None:
        col = hex_of(solid)
        if col:
            return f"background:{rgba(col, alpha_of(solid))};"
    grad = spPr.find("a:gradFill", NS)
    if grad is not None:
        stops = []
        for gs in grad.findall("a:gsLst/a:gs", NS):
            col = hex_of(gs) or "#FFFFFF"
            pos = int(gs.get("pos", "0")) / 1000.0
            stops.append(f"{rgba(col, alpha_of(gs))} {round(pos, 1)}%")
        ang = grad.find("a:lin", NS)
        deg = 90
        if ang is not None:
            # DrawingML measures clockwise from 3 o'clock in 1/60000 degrees;
            # CSS measures clockwise from 12 o'clock
            deg = (int(ang.get("ang", "0")) / 60000.0 + 90) % 360
        if stops:
            return f"background:linear-gradient({round(deg)}deg,{','.join(stops)});"
    return ""


def line_css(shape) -> str:
    ln = shape._element.find(".//a:ln", NS)
    if ln is None:
        return ""
    if ln.find("a:noFill", NS) is not None:
        return ""
    solid = ln.find("a:solidFill", NS)
    if solid is None:
        return ""
    col = hex_of(solid)
    if not col:
        return ""
    w = float(ln.get("w", "9525")) / 914400.0 * PX
    return f"border:{max(round(w, 2), 1)}px solid {rgba(col, alpha_of(solid))};"


def radius_css(shape) -> str:
    try:
        name = str(shape.shape_type)
    except Exception:
        return ""
    if "ROUNDED" in name.upper():
        return "border-radius:10px;"
    return ""


def run_css(run, para) -> str:
    f = run.font
    out = []
    size = f.size or para.font.size
    if size is not None:
        out.append(f"font-size:{round(size.pt, 1)}pt")
    name = f.name or para.font.name
    if name:
        out.append(f"font-family:'{name}',serif")
    bold = f.bold if f.bold is not None else para.font.bold
    if bold:
        out.append("font-weight:700")
    if f.italic:
        out.append("font-style:italic")
    if f.underline:
        out.append("text-decoration:underline")
    try:
        if f.color is not None and f.color.type is not None and f.color.rgb is not None:
            out.append(f"color:#{f.color.rgb}")
    except Exception:
        pass
    return ";".join(out)


def text_html(shape) -> str:
    tf = shape.text_frame
    parts = []
    for para in tf.paragraphs:
        styles = []
        if para.alignment is not None:
            styles.append(f"text-align:{ALIGN.get(int(para.alignment), 'left')}")
        if para.space_before is not None:
            styles.append(f"margin-top:{round(para.space_before.pt, 1)}pt")
        if para.space_after is not None:
            styles.append(f"margin-bottom:{round(para.space_after.pt, 1)}pt")
        if para.line_spacing is not None and not isinstance(para.line_spacing, float):
            styles.append(f"line-height:{round(para.line_spacing.pt, 1)}pt")
        elif isinstance(para.line_spacing, float):
            styles.append(f"line-height:{para.line_spacing}")
        runs = "".join(
            f'<span style="{run_css(r, para)}">{html.escape(r.text).replace(chr(10), "<br>")}</span>'
            for r in para.runs) or "<br>"
        parts.append(f'<p style="{";".join(styles)}">{runs}</p>')
    return "".join(parts)


def picture_html(shape) -> str:
    img = shape.image
    b64 = base64.b64encode(img.blob).decode("ascii")
    return (f'<img src="data:{img.content_type};base64,{b64}" '
            f'alt="figure" draggable="false">')


def table_html(shape) -> str:
    rows = []
    for r in shape.table.rows:
        cells = []
        for c in r.cells:
            style = []
            # only the cell properties carry the cell fill; searching the whole
            # cell would pick up a run's font colour and paint the cell with it
            solid = c._tc.find("./a:tcPr/a:solidFill", NS)
            if solid is not None:
                col = hex_of(solid)
                if col:
                    style.append(f"background:{rgba(col, alpha_of(solid))}")
            body = []
            for para in c.text_frame.paragraphs:
                al = ALIGN.get(int(para.alignment), "left") if para.alignment is not None else "left"
                runs = "".join(
                    f'<span style="{run_css(rr, para)}">{html.escape(rr.text)}</span>'
                    for rr in para.runs) or "&nbsp;"
                body.append(f'<p style="text-align:{al}">{runs}</p>')
            cells.append(f'<td style="{";".join(style)}">{"".join(body)}</td>')
        rows.append("<tr>" + "".join(cells) + "</tr>")
    return f'<table class="tbl">{"".join(rows)}</table>'


def shape_div(shape, dx=0.0, dy=0.0, sx=1.0, sy=1.0) -> str:
    """Render one shape. dx/dy/sx/sy carry the parent group transform."""
    if shape.left is None or shape.top is None:
        return ""
    left = dx + px(shape.left) * sx
    top = dy + px(shape.top) * sy
    w = px(shape.width) * sx
    h = px(shape.height) * sy
    style = [f"left:{round(left, 2)}px", f"top:{round(top, 2)}px",
             f"width:{round(w, 2)}px", f"height:{round(h, 2)}px"]
    rot = getattr(shape, "rotation", 0) or 0
    if rot:
        style.append(f"transform:rotate({rot}deg)")

    kind = "shape"
    inner = ""
    if shape.shape_type == 13 or shape.__class__.__name__ == "Picture":
        kind = "pic"
        inner = picture_html(shape)
    elif getattr(shape, "has_table", False) and shape.has_table:
        kind = "tbl-wrap"
        inner = table_html(shape)
    else:
        css = fill_css(shape) + line_css(shape) + radius_css(shape)
        if css:
            style.append(css.rstrip(";"))
        if shape.has_text_frame and shape.text_frame.text.strip():
            tf = shape.text_frame
            anchor = ANCHOR.get(int(tf.vertical_anchor), "flex-start") \
                if tf.vertical_anchor is not None else "flex-start"
            style.append(f"justify-content:{anchor}")
            if tf.word_wrap is False:
                style.append("white-space:nowrap")
            pad = []
            for side, attr in (("top", "margin_top"), ("right", "margin_right"),
                               ("bottom", "margin_bottom"), ("left", "margin_left")):
                v = getattr(tf, attr, None)
                pad.append(f"{round(px(v), 1)}px" if v is not None else "0px")
            style.append("padding:" + " ".join(pad))
            inner = f'<div class="txt" contenteditable="true">{text_html(shape)}</div>'

    return (f'<div class="obj {kind}" data-kind="{kind}" '
            f'style="{";".join(style)}">{inner}</div>')


def walk(shapes, out, dx=0.0, dy=0.0, sx=1.0, sy=1.0):
    for sh in shapes:
        if sh.shape_type == 6:  # group: map child coordinate space onto the slide
            grp = sh._element.find(".//a:xfrm", NS)
            csx, csy = sx, sy
            ndx, ndy = dx, dy
            if grp is not None:
                off = grp.find("a:off", NS)
                ext = grp.find("a:ext", NS)
                choff = grp.find("a:chOff", NS)
                chext = grp.find("a:chExt", NS)
                if None not in (off, ext, choff, chext):
                    csx = sx * (int(ext.get("cx")) / max(int(chext.get("cx")), 1))
                    csy = sy * (int(ext.get("cy")) / max(int(chext.get("cy")), 1))
                    ndx = dx + px(off.get("x")) * sx - px(choff.get("x")) * csx
                    ndy = dy + px(off.get("y")) * sy - px(choff.get("y")) * csy
            walk(sh.shapes, out, ndx, ndy, csx, csy)
            continue
        out.append(shape_div(sh, dx, dy, sx, sy))


PAGE_CSS = """
:root { --w: 1280px; --h: 720px; }
* { box-sizing: border-box; }
body { margin: 0; background: #4A5750; font-family: 'Times New Roman', serif; }
#bar { position: fixed; inset: 0 0 auto 0; height: 46px; z-index: 50;
       background: #2F5248; color: #fff; display: flex; align-items: center;
       gap: 14px; padding: 0 18px; font: 13px/1 'Segoe UI', sans-serif; }
#bar b { font: 600 13px/1 'Segoe UI', sans-serif; letter-spacing: .06em; }
#bar button { font: 12px/1 'Segoe UI', sans-serif; padding: 7px 12px; border: 0;
              border-radius: 5px; background: #E8F0E9; color: #2F5248; cursor: pointer; }
#bar button:hover { background: #fff; }
#bar span.hint { opacity: .75; margin-left: auto; }
#deck { padding: 70px 0 60px; }
.slide { position: relative; width: var(--w); height: var(--h); margin: 0 auto 26px;
         background: #fff; overflow: hidden; box-shadow: 0 8px 26px rgba(0,0,0,.35); }
.slide .num { position: absolute; right: -46px; top: 6px; color: #cfd8d2;
              font: 12px/1 'Segoe UI', sans-serif; }
.obj { position: absolute; display: flex; flex-direction: column; }
.obj .txt { width: 100%; outline: none; }
.obj .txt p { margin: 0; }
.pic img { width: 100%; height: 100%; object-fit: contain; }
.tbl { width: 100%; height: 100%; border-collapse: collapse; }
.tbl td { border: 1px solid #D5DED6; padding: 3px 6px; vertical-align: middle; }
.tbl td p { margin: 0; }
body.edit .obj:hover { outline: 1px dashed rgba(47,82,72,.55); }
body.edit .obj.sel { outline: 2px solid #C8102E; }
body.edit .obj.sel::after { content: ''; position: absolute; right: -6px; bottom: -6px;
       width: 12px; height: 12px; background: #C8102E; cursor: nwse-resize; }
@media print {
  body { background: #fff; } #bar { display: none; } #deck { padding: 0; }
  .slide { margin: 0; box-shadow: none; page-break-after: always; }
  .slide .num { display: none; }
}
"""

SCRIPT = """
const body = document.body;
let editMode = false, sel = null, drag = null;

function setEdit(on) {
  editMode = on;
  body.classList.toggle('edit', on);
  document.getElementById('btnEdit').textContent =
    on ? 'Layout editing: ON' : 'Layout editing: OFF';
  if (!on && sel) { sel.classList.remove('sel'); sel = null; }
}

document.getElementById('btnEdit').onclick = () => setEdit(!editMode);

document.addEventListener('mousedown', e => {
  if (!editMode) return;
  const obj = e.target.closest('.obj');
  if (!obj) { if (sel) sel.classList.remove('sel'); sel = null; return; }
  if (e.target.isContentEditable && !e.altKey) return;
  if (sel) sel.classList.remove('sel');
  sel = obj; sel.classList.add('sel');
  const r = obj.getBoundingClientRect();
  const corner = (e.clientX > r.right - 14) && (e.clientY > r.bottom - 14);
  drag = { obj, corner, x: e.clientX, y: e.clientY,
           left: parseFloat(obj.style.left), top: parseFloat(obj.style.top),
           w: parseFloat(obj.style.width), h: parseFloat(obj.style.height) };
  e.preventDefault();
});

document.addEventListener('mousemove', e => {
  if (!drag) return;
  const dx = e.clientX - drag.x, dy = e.clientY - drag.y;
  if (drag.corner) {
    drag.obj.style.width = Math.max(12, drag.w + dx) + 'px';
    drag.obj.style.height = Math.max(12, drag.h + dy) + 'px';
  } else {
    drag.obj.style.left = (drag.left + dx) + 'px';
    drag.obj.style.top = (drag.top + dy) + 'px';
  }
});

document.addEventListener('mouseup', () => { drag = null; });

document.addEventListener('keydown', e => {
  if (!editMode || !sel) return;
  if (e.key === 'Delete' || e.key === 'Backspace') {
    if (document.activeElement && document.activeElement.isContentEditable) return;
    sel.remove(); sel = null; e.preventDefault();
    return;
  }
  // while the caret is in a text box, arrows belong to the caret
  if (document.activeElement && document.activeElement.isContentEditable) return;
  const step = e.shiftKey ? 10 : 1;
  const move = { ArrowLeft: [-step, 0], ArrowRight: [step, 0],
                 ArrowUp: [0, -step], ArrowDown: [0, step] }[e.key];
  if (move) {
    // fall back on the measured box, in case the shape has no inline geometry
    const box = sel.getBoundingClientRect();
    const par = sel.offsetParent.getBoundingClientRect();
    const l = parseFloat(sel.style.left);
    const t = parseFloat(sel.style.top);
    sel.style.left = ((isNaN(l) ? box.left - par.left : l) + move[0]) + 'px';
    sel.style.top = ((isNaN(t) ? box.top - par.top : t) + move[1]) + 'px';
    e.preventDefault();
  }
}, true);   // capture, so a focused text box cannot swallow the arrow key

// swap any figure for a local image file
document.querySelectorAll('.pic img').forEach(img => {
  img.addEventListener('dblclick', () => {
    const inp = document.createElement('input');
    inp.type = 'file'; inp.accept = 'image/*';
    inp.onchange = () => {
      const f = inp.files[0]; if (!f) return;
      const fr = new FileReader();
      fr.onload = () => { img.src = fr.result; };
      fr.readAsDataURL(f);
    };
    inp.click();
  });
});

document.getElementById('btnSave').onclick = () => {
  const clone = document.documentElement.cloneNode(true);
  clone.querySelectorAll('.obj.sel').forEach(o => o.classList.remove('sel'));
  clone.querySelector('body').classList.remove('edit');
  const blob = new Blob(['<!DOCTYPE html>\\n' + clone.outerHTML],
                        { type: 'text/html;charset=utf-8' });
  const a = document.createElement('a');
  a.href = URL.createObjectURL(blob);
  a.download = 'CA6002_Group30_Final_Presentation_edited.html';
  a.click();
};

document.getElementById('btnPrint').onclick = () => window.print();
"""


def main() -> None:
    prs = Presentation(DECK)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        out = []
        walk(slide.shapes, out)
        slides.append(f'<section class="slide" data-slide="{i}">'
                      f'<div class="num">{i}</div>{"".join(out)}</section>')

    w = round(prs.slide_width / 914400 * PX)
    h = round(prs.slide_height / 914400 * PX)
    doc = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>CA6002 Group 30 - Investigating Country Carbon Emission Patterns</title>
<style>{PAGE_CSS}
:root {{ --w: {w}px; --h: {h}px; }}
</style>
</head>
<body>
<div id="bar">
  <b>CA6002 · GROUP 30 · EDITABLE HTML DECK</b>
  <button id="btnEdit">Layout editing: OFF</button>
  <button id="btnSave">Save as HTML</button>
  <button id="btnPrint">Print / PDF</button>
  <span class="hint">Click any text to edit it · double-click a figure to replace it ·
  turn on layout editing to drag, resize (corner) or delete shapes</span>
</div>
<div id="deck">
{chr(10).join(slides)}
</div>
<script>{SCRIPT}</script>
</body>
</html>
"""
    OUT.write_text(doc, encoding="utf-8")
    mb = OUT.stat().st_size / 1e6
    print(f"saved: {OUT}  slides: {len(slides)}  size: {mb:.1f} MB")


if __name__ == "__main__":
    main()
